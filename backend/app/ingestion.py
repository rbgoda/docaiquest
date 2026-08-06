"""Ingestion pipeline.

Three pure-ish stages composed into `ingest_document`:

1. **Parse** — pull text + page numbers out of the stored PDF via PyMuPDF.
2. **Chunk** — sliding ~1000-char windows with overlap, preserving page.
3. **Embed** — vectorize each chunk through the configured backend.

The composed function persists the chunks + flips the document's
`ingestion_status` (`pending` → `processing` → `ready` / `failed`). It's
called from the Arq worker; nothing in the request path runs it.
"""

from __future__ import annotations

import hashlib
import io
import logging
import os
import re
from dataclasses import dataclass

import fitz  # PyMuPDF
from sqlalchemy import delete, select
from sqlalchemy.orm import Session

from app import ingestion_vision, storage
from app.config import get_settings
from app.db import set_current_tenant
from app.embeddings import embed
from app.feature_flags import is_enabled, get_int
from app.entities import extract_entities
from app.orm import Document, DocumentChunk, Entity

log = logging.getLogger("docaiq.ingestion")


# C0 control bytes that PostgreSQL TEXT rejects (NUL, 0x00) or that are pure
# extraction noise — keep the legitimate whitespace controls (\t \n \r). Some
# PDFs (complex-script CID fonts, malformed encodings) and some OCR output yield
# extracted text containing NUL, which otherwise crashes the chunk INSERT with
# `psycopg.DataError: ... cannot contain NUL (0x00) bytes` and fails ingestion.
_CTRL_BYTES_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def _sanitize_text(s: str | None) -> str:
    """Strip NUL + other C0 control bytes (keeps tab/newline/carriage-return) so
    extracted text is safe to persist in Postgres TEXT and clean for the LLM."""
    if not s:
        return s or ""
    return _CTRL_BYTES_RE.sub("", s)


def _decode_csv(raw: bytes) -> str:
    """Decode CSV bytes into text.

    T3.1 · chardet detection first (high-confidence cases pass through
    cleanly), then a sequential try-list for utf/cp/latin variants.
    Hard-fails with `ValueError` only when EVERY codec corrupts — the
    caller (ingest_document) wraps in an ingestion_error so the user
    sees a clear message instead of garbled output downstream."""
    # 1. chardet · trust high-confidence detections (>0.65). Latin-1 always
    #    "succeeds" so we don't use it for confidence, just as a fallback.
    try:
        import chardet
        result = chardet.detect(raw[:65536])
        enc = (result or {}).get("encoding")
        conf = (result or {}).get("confidence", 0.0) or 0.0
        if enc and conf >= 0.65 and enc.lower() != "ascii":
            try:
                return raw.decode(enc)
            except (UnicodeDecodeError, LookupError):
                pass
    except ImportError:
        pass  # chardet missing — fall through to try-list
    # 2. Ordered try-list of common bank/CC export encodings.
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1", "utf-16"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    # 3. Last resort · lossy utf-8. Log so operators can see this happened.
    log.warning("CSV: all encodings failed cleanly, falling back to lossy utf-8")
    return raw.decode("utf-8", errors="replace")


@dataclass(frozen=True)
class Chunk:
    text: str
    page: int
    char_start: int
    char_end: int
    # P9.5 · 'text' (sliding-window text) or 'table' (Markdown table from
    # pdfplumber). Defaults to 'text' so every existing call site is unchanged.
    kind: str = "text"
    # Phase 5 · forward bbox provenance from the IR (union of the composing blocks'
    # bboxes). None on the flat path; the search_for locate fills bboxes there.
    bbox: dict | None = None
    # Camelot spatial bbox (x0, y0, x1, y1) in PDF points, bottom-left origin.
    # Used to find matching line_map lines for precise per-chunk highlighting
    # instead of the full-page fallback.  None for non-table or pdfplumber chunks.
    table_bbox: tuple[float, float, float, float] | None = None
    # Composing document-level block IDs as string tuples (e.g. ("b_0000", "b_0001")).
    # Empty tuple on the flat path (chunk_pages) where there's no block registry.
    # Converted to a JSONB list when written to DocumentChunk.block_ids.
    block_ids: tuple[str, ...] = ()


def _blocks_from_docling_items(docling_doc) -> list:
    """Build IR Block objects directly from a Docling DoclingDocument, preserving
    per-item bbox provenance from `item.prov[0].bbox` + `item.prov[0].page_no`.

    This bypasses the lossy ``export_to_markdown()`` → ``blocks_from_markdown()``
    round-trip so every block carries its real page + bounding box.  Downstream
    ``chunking.chunk_blocks`` unions these into chunk bboxes automatically —
    no word-matching / spatial fallback needed.

    DocItemLabel → BlockKind mapping:
    - TITLE, SECTION_HEADER  → HEADING (level from item.level)
    - TEXT, PARAGRAPH, CAPTION, FOOTNOTE, PAGE_HEADER, PAGE_FOOTER,
      CODE, FORMULA, FORM, KEY_VALUE_REGION, CHECKBOX_*, FIELD_*,
      HANDWRITTEN_TEXT, REFERENCE, MARKER, EMPTY_VALUE → PARAGRAPH
    - LIST_ITEM → LIST_ITEM
    - TABLE → TABLE (grid reconstructed from item.data.table_cells)
    - PICTURE, CHART → FIGURE
    """
    from app.document_model import Block, BlockKind, BBox

    # Resolve page dimensions so we can compute percentage-based bboxes.
    _page_dims: dict[int, tuple[float, float]] = {}
    if hasattr(docling_doc, "pages") and docling_doc.pages:
        for pno, page_info in docling_doc.pages.items():
            sz = getattr(page_info, "size", None)
            if sz is not None:
                _page_dims[pno] = (getattr(sz, "width", 0) or 0,
                                   getattr(sz, "height", 0) or 0)

    def _make_bbox(prov, page: int) -> "BBox | None":
        """Convert a Docling ProvenanceItem.bbox → our BBox (TOPLEFT origin)."""
        if prov is None:
            return None
        bb = getattr(prov, "bbox", None)
        if bb is None:
            return None
        pw, ph = _page_dims.get(page, (0, 0))
        # Docling bbox: l, t, r, b.  coord_origin is TOPLEFT or BOTTOMLEFT.
        origin = str(getattr(bb, "coord_origin", "TOPLEFT") or "TOPLEFT")
        x0, x1 = bb.l, bb.r
        if "BOTTOM" in origin.upper() and ph > 0:
            # In BOTTOMLEFT, t (top) > b (bottom). TOPLEFT flips: y0=top<bottom=y1.
            y0, y1 = ph - bb.t, ph - bb.b
        else:
            y0, y1 = bb.t, bb.b
        return BBox(x0=x0, y0=y0, x1=x1, y1=y1, page_w=pw, page_h=ph, page=page)

    TEXT_LABELS = {
        "text", "paragraph", "caption", "footnote", "page_header", "page_footer",
        "code", "formula", "reference", "marker", "handwritten_text",
        "form", "key_value_region", "checkbox_selected", "checkbox_unselected",
        "field_key", "field_value", "field_heading", "field_hint", "field_item",
        "field_region", "empty_value",
    }

    blocks: list[Block] = []
    for item, _level in docling_doc.iterate_items():
        label = str(getattr(item, "label", "text") or "text")
        if hasattr(label, "value"):
            label = label.value

        prov = item.prov[0] if (item.prov and len(item.prov) > 0) else None
        page = prov.page_no if prov else 1
        bbox = _make_bbox(prov, page)

        if label in ("title", "section_header"):
            lvl = getattr(item, "level", 1) or 1
            text = getattr(item, "text", "") or ""
            blocks.append(Block(kind=BlockKind.HEADING, page=page, text=text,
                                level=lvl, bbox=bbox))

        elif label == "list_item":
            text = getattr(item, "text", "") or ""
            blocks.append(Block(kind=BlockKind.LIST_ITEM, page=page, text=text,
                                bbox=bbox))

        elif label == "table":
            data = getattr(item, "data", None)
            if data is not None and hasattr(data, "table_cells") and data.table_cells:
                nr = getattr(data, "num_rows", 0) or 0
                nc = getattr(data, "num_cols", 0) or 0
                grid = [["" for _ in range(nc)] for _ in range(nr)]
                for cell in data.table_cells:
                    r, c = getattr(cell, "row", 0), getattr(cell, "col", 0)
                    if 0 <= r < nr and 0 <= c < nc:
                        grid[r][c] = getattr(cell, "text", "") or ""
                has_header = nr > 0 and any(
                    getattr(c, "column_header", False)
                    for c in data.table_cells
                    if getattr(c, "row", 0) == 0
                )
                blocks.append(Block(kind=BlockKind.TABLE, page=page, rows=grid,
                                    has_header=has_header, bbox=bbox))

        elif label in ("picture", "chart"):
            caption = ""
            captions = getattr(item, "captions", None) or []
            if captions:
                caption = getattr(captions[0], "text", "") or ""
            blocks.append(Block(kind=BlockKind.FIGURE, page=page, caption=caption,
                                bbox=bbox))

        elif label in TEXT_LABELS:
            text = getattr(item, "text", "") or ""
            if text.strip():
                blocks.append(Block(kind=BlockKind.PARAGRAPH, page=page,
                                    text=text, bbox=bbox))

        else:
            text = getattr(item, "text", "") or ""
            if text.strip():
                blocks.append(Block(kind=BlockKind.PARAGRAPH, page=page,
                                    text=text, bbox=bbox))

    return blocks


def _fill_docling_gaps(
    blocks: list, pdf_bytes: bytes,
) -> list:
    """Fill text that Docling missed by comparing with PyMuPDF's extraction.

    Docling occasionally classifies text regions as pictures (especially in
    headers on later pages of multi-page documents), losing critical content
    like amounts, IDs, and dates.  PyMuPDF's ``page.get_text("text")`` is a
    reliable fallback — it reads the embedded text stream directly.

    Returns a new block list (may be the same object if no gaps found)."""
    import fitz

    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            pymu_pages = [(i + 1, page.get_text("text"))
                          for i, page in enumerate(doc)]
    except Exception:  # noqa: BLE001
        return blocks

    # Build a set of normalized text snippets already covered by Docling blocks,
    # per page.  We compare short substrings (≥8 chars after normalization) so a
    # single missed word doesn't trigger a duplicate block.
    _MIN_SUBSTR = 20
    docling_text: dict[int, set[str]] = {}
    for b in blocks:
        t = (getattr(b, "text", "") or "").strip()
        if not t:
            continue
        pg = getattr(b, "page", 1) or 1
        docling_text.setdefault(pg, set())
        # Add the full text and all substrings ≥ _MIN_SUBSTR chars
        docling_text[pg].add(t.lower())
        for j in range(len(t) - _MIN_SUBSTR + 1):
            docling_text[pg].add(t[j:j + _MIN_SUBSTR].lower())

    # Per-page gap check
    from app.document_model import Block, BlockKind
    added = 0
    for page_no, pymu_text in pymu_pages:
        if not pymu_text or not pymu_text.strip():
            continue
        covered = docling_text.get(page_no, set())
        # Split PyMuPDF text into lines; check each line for coverage
        for line in pymu_text.split("\n"):
            line_s = line.strip()
            if len(line_s) < _MIN_SUBSTR:
                continue
            # Check if any 12-char substring of this line is already covered
            line_lower = line_s.lower()
            found = any(
                line_lower[j:j + _MIN_SUBSTR] in covered
                for j in range(len(line_lower) - _MIN_SUBSTR + 1)
            )
            if not found:
                # PyMuPDF text not in Docling blocks — add as a gap-fill block
                blocks.append(Block(
                    kind=BlockKind.PARAGRAPH, page=page_no, text=line_s,
                    # No bbox from Docling; the line_map will cover it later
                ))
                added += 1
    if added:
        # Re-sort by page so gap-fill blocks land in reading order instead
        # of being appended to the end of the list.
        blocks.sort(key=lambda b: getattr(b, "page", 1) or 1)
        log.info("docling: gap-filled %d missing text line(s) from PyMuPDF", added)
    return blocks


def parse_pdf_docling(
    pdf_bytes: bytes,
) -> tuple[list[tuple[int, str]], object] | None:
    """G13 · Parse PDF via Docling (MIT, layout-aware). Returns (pages, ir_document)
    with proper multi-column reading order, table structure, and figure detection.
    Returns None if Docling is not installed or fails — caller falls back to PyMuPDF.

    Docling handles natively:
      · Multi-column layouts (no CV vocabulary gate)
      · Complex tables with merged cells
      · Figure/caption binding
      · Heading hierarchy
      · Reading order reconstruction
    """
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        return None
    try:
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
            tf.write(pdf_bytes)
            tmp = tf.name
        try:
            converter = DocumentConverter()
            result = converter.convert(tmp)
            # export_to_markdown() is only needed for the fallback path below;
            # the happy path builds blocks directly from Docling items.
        finally:
            os.unlink(tmp)

        if not result.document:
            return None

        # Build IR blocks directly from Docling items — preserves per-block
        # bbox provenance so chunk bboxes are exact without word-matching.
        from app.document_model import Document
        try:
            blocks = _blocks_from_docling_items(result.document)
            if not blocks:
                return None
        except Exception as e:  # noqa: BLE001 — fall back to markdown path
            log.info("docling: direct-item parse failed (%s); falling back to markdown", e)
            md = result.document.export_to_markdown()
            if not md or not md.strip():
                return None
            from app.markdown_ir import blocks_from_markdown
            blocks = blocks_from_markdown(md)

        # Fill gaps where Docling classified text as pictures (common on
        # later pages of multi-page PDFs — amounts, IDs, dates go missing).
        blocks = _fill_docling_gaps(blocks, pdf_bytes)

        ir = Document(blocks=blocks)
        pages = ir.to_pages()
        _with_bbox = sum(1 for b in blocks if b.bbox is not None)
        log.info("docling: parsed %d pages, %d blocks (%d with bbox)",
                 len(pages), len(blocks), _with_bbox)
        return pages, ir
    except Exception as e:
        log.info("docling: parse failed, falling back to PyMuPDF: %s", e)
        return None


def parse_pdf(pdf_bytes: bytes) -> list[tuple[int, str]]:
    """Return [(page_number_1based, text), ...]. Empty pages are kept so the
    chunk index aligns with the document's natural page order."""
    pages: list[tuple[int, str]] = []
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc, start=1):
            pages.append((i, page.get_text("text")))
    return pages


# ── Phase 1 · structured PDF parse → Document Model (IR) ────────────────────
# Bold flag bit in a PyMuPDF span's `flags`. Used only as a weak heading signal.
_FITZ_BOLD = 1 << 4  # 16

# ── Multi-column reading-order reconstruction ───────────────────────────────
# PyMuPDF groups text into blocks that can MERGE across a two-column layout — a
# résumé's left "EDUCATION" column and right "PROJECTS" column land in ONE block
# on the same y-band — so linearised text interleaves unrelated sections and
# attaches the wrong dates/values to each entry (e.g. an HSC row picks up a
# project's year, swapping HSC↔SSC in chat/extraction). We rebuild reading order
# at the WORD level: detect the vertical gutter, split full-width banners into
# bands, and emit each column top→bottom, left→right. Conservative — only fires
# on pages that are unambiguously two-column; everything else keeps the native
# block parse, and any error falls back silently.
_COL_ROW_Y_TOL = 5.0        # words within this y-gap belong to the same visual row
_COL_FULLW_FRAC = 0.62      # a row this wide (of page width) spanning the gutter = full-width band
_COL_GUTTER_LO, _COL_GUTTER_HI = 0.30, 0.70   # search the gutter only in the middle of the page


def _col_cluster_rows(words):
    """Group PyMuPDF words [(x0,y0,x1,y1,text,…)] into visual rows by y proximity."""
    rows: list[dict] = []
    for w in sorted(words, key=lambda w: (round(w[1]), w[0])):
        cy = (w[1] + w[3]) / 2.0
        for r in rows:
            if abs(cy - r["cy"]) <= _COL_ROW_Y_TOL:
                r["w"].append(w); r["ys"].append(cy)
                r["cy"] = sum(r["ys"]) / len(r["ys"])
                break
        else:
            rows.append({"cy": cy, "ys": [cy], "w": [w]})
    rows.sort(key=lambda r: r["cy"])
    return rows


def _col_detect_gutter(rows, page_w: float):
    """Return the gutter x if the page is cleanly two-column, else None. A gutter
    is the x that maximises rows split cleanly (words on both sides, none crossing)
    while few rows cross it — single-column pages never satisfy both."""
    if len(rows) < 8:
        return None
    best = None
    for k in range(int(_COL_GUTTER_LO * 100), int(_COL_GUTTER_HI * 100) + 1):
        X = page_w * k / 100.0
        split = crossing = 0
        for r in rows:
            if any(w[0] < X < w[2] for w in r["w"]):
                crossing += 1
            elif any(w[2] <= X for w in r["w"]) and any(w[0] >= X for w in r["w"]):
                split += 1
        if best is None or split > best[1] or (split == best[1] and crossing < best[2]):
            best = (X, split, crossing)
    X, split, crossing = best
    if split >= max(4, int(0.35 * len(rows))) and crossing <= 0.25 * max(1, split):
        return X
    return None


# Column reconstruction is a RÉSUMÉ/CV-specific fix — that's the layout where two
# columns are genuinely independent titled sections. A document-processing corpus is
# otherwise dominated by tables / statements / forms / lab & medical reports, where
# column-major reading SCRAMBLES row-linked data. So we gate reconstruction on résumé
# CONTENT: a CV page carries a recognisable vocabulary of section headers that those
# other document types do not. Require several distinct hits before reconstructing.
_RESUME_SECTIONS = frozenset({
    "objective", "career objective", "summary", "professional summary", "profile",
    "education", "academic", "qualification", "qualifications",
    "experience", "work experience", "employment", "work history",
    "skills", "technical skills", "soft skills", "core competencies", "competencies",
    "projects", "certifications", "certification", "achievements", "accomplishments",
    "languages", "languages known", "interests", "hobbies", "core interest",
    "references", "awards", "publications", "activities", "participations",
    "extracurricular", "declaration", "strengths", "personal details", "contact",
})


def _resume_section_hits(rows, gutter: float) -> int:
    """Count DISTINCT résumé section headers among the page's rows — the content
    signature that separates a CV from a table/statement/form/report. Each row is
    split at the gutter so column-aligned headers (EDUCATION│PROJECTS land in one
    PyMuPDF row) are both counted."""
    seen: set[str] = set()
    for r in rows:
        for side in ([w for w in r["w"] if w[2] <= gutter], [w for w in r["w"] if w[0] >= gutter]):
            txt = " ".join(w[4] for w in sorted(side, key=lambda w: w[0])).strip().lower().rstrip(":").strip()
            if txt in _RESUME_SECTIONS:
                seen.add(txt)
    return len(seen)


def _col_is_heading(text: str) -> bool:
    """An all-caps short alpha label (EDUCATION, TECHNICAL SKILLS, HSC) reads as a
    section heading. Digit/long/mixed-case lines stay paragraphs."""
    t = text.strip()
    if not t or len(t.split()) > 6 or any(ch.isdigit() for ch in t):
        return False
    return sum(ch.isalpha() for ch in t) >= 2 and t.isupper()


def _page_to_blocks_columnar(page, pno: int):
    """Word-level column reconstruction → IR blocks, or None when the page is not
    cleanly multi-column (caller then falls back to the native block parser)."""
    from app.document_model import BBox, Block, BlockKind, SOURCE_NATIVE
    page_w = float(page.rect.width or 0.0)
    page_h = float(page.rect.height or 0.0)
    if not page_w:
        return None
    words = page.get_text("words")
    if len(words) < 12:
        return None
    rows = _col_cluster_rows(words)
    gutter = _col_detect_gutter(rows, page_w)
    if gutter is None:
        return None
    # Decisive gate: only reconstruct pages whose content reads as a résumé/CV — the
    # one layout where independent columns are the norm. Tables, statements, forms,
    # lab & medical reports lack this section vocabulary, so they keep the native
    # (row-major) parse and their row-linked data is never scrambled.
    if _resume_section_hits(rows, gutter) < 3:
        return None
    fullw = _COL_FULLW_FRAC * page_w

    # Segment into horizontal bands delimited by full-width rows (banners / rules).
    bands: list[tuple[str, list]] = []
    cur: list = []
    for r in rows:
        xs = [w[0] for w in r["w"]] + [w[2] for w in r["w"]]
        is_full = any(w[0] < gutter < w[2] for w in r["w"]) and (max(xs) - min(xs)) >= fullw
        if is_full:
            if cur:
                bands.append(("cols", cur)); cur = []
            bands.append(("full", [r]))
        else:
            cur.append(r)
    if cur:
        bands.append(("cols", cur))

    # Guard against row-linked TABLES (a lab report's test-name column paired with a
    # value column): column-major reading would scramble their row pairings. Only
    # genuine INDEPENDENT columns — a résumé's left and right are separate titled
    # sections — carry section headings in BOTH columns; a data table does not.
    # Require ≥2 all-caps heading rows on each side, else keep the native row-major
    # parse (no worse than today for non-résumé two-column pages).
    lh = rh = 0
    for _kind, brows in bands:
        if _kind != "cols":
            continue
        for r in brows:
            lw = sorted([w for w in r["w"] if w[2] <= gutter], key=lambda w: w[0])
            rw = sorted([w for w in r["w"] if w[0] >= gutter], key=lambda w: w[0])
            if lw and _col_is_heading(" ".join(w[4] for w in lw)):
                lh += 1
            if rw and _col_is_heading(" ".join(w[4] for w in rw)):
                rh += 1
    if lh < 2 or rh < 2:
        return None

    # Second guard: a row-linked DATA TABLE has ≥3 columns (Test | Result | Range),
    # so the right side of the primary gutter splits AGAIN — two INDEPENDENT columns
    # (a résumé) do not. And its cells carry table-header words. If the right column
    # subdivides into its own gutter, or ≥2 of its rows read as table headers, treat
    # the page as tabular and keep the native parse. (Lab/financial reports repeatedly
    # defeated the heading count alone; this cleanly separates them from résumés.)
    right_sub = [{"w": [w for w in r["w"] if w[0] >= gutter]} for r in rows]
    right_sub = [r for r in right_sub if r["w"]]
    if _col_detect_gutter(right_sub, page_w) is not None:
        return None
    _HDR_WORDS = {"result", "results", "unit", "units", "ref", "reference", "range", "test", "value"}
    hdr_hits = 0
    for r in rows:
        low = " ".join(w[4] for w in r["w"] if w[0] >= gutter).lower()
        if _HDR_WORDS & set(low.split()):
            hdr_hits += 1
    if hdr_hits >= 2:
        return None

    blocks: list = []

    def _mk(text: str, ws: list, heading: bool):
        text = text.strip()
        if not text:
            return
        bb = None
        if ws and page_h:
            bb = BBox(float(min(w[0] for w in ws)), float(min(w[1] for w in ws)),
                      float(max(w[2] for w in ws)), float(max(w[3] for w in ws)),
                      page_w, page_h, pno)
        kind, level = (BlockKind.HEADING, 2) if heading else (BlockKind.PARAGRAPH, 0)
        blocks.append(Block(kind=kind, page=pno, text=text, level=level,
                            bbox=bb, source=SOURCE_NATIVE))

    def _emit_col(colrows: list):
        para_txt: list[str] = []
        para_ws: list = []
        for r in colrows:
            ws = sorted(r["w"], key=lambda w: w[0])
            txt = " ".join(w[4] for w in ws)
            if not txt.strip():
                continue
            if _col_is_heading(txt):
                if para_txt:
                    _mk("\n".join(para_txt), para_ws, False); para_txt, para_ws = [], []
                _mk(txt, ws, True)
            else:
                para_txt.append(txt); para_ws.extend(ws)
        if para_txt:
            _mk("\n".join(para_txt), para_ws, False)

    for kind, brows in bands:
        if kind == "full":
            ws = sorted(brows[0]["w"], key=lambda w: w[0])
            _mk(" ".join(w[4] for w in ws), ws, False)
            continue
        left, right = [], []
        for r in brows:
            l = [w for w in r["w"] if w[2] <= gutter]
            rr = [w for w in r["w"] if w[0] >= gutter]
            for w in r["w"]:                     # stray word crossing the gutter → its center side
                if w[0] < gutter < w[2]:
                    (l if (w[0] + w[2]) / 2 < gutter else rr).append(w)
            if l:
                left.append({"cy": r["cy"], "w": l})
            if rr:
                right.append({"cy": r["cy"], "w": rr})
        _emit_col(sorted(left, key=lambda r: r["cy"]))
        _emit_col(sorted(right, key=lambda r: r["cy"]))
    return blocks or None


def _page_to_blocks(page, pno: int):
    """Turn one text-layer PDF page into IR blocks via PyMuPDF `get_text("dict")`.

    On a cleanly two-column page, word-level column reconstruction runs first so
    the left column is read fully before the right (native blocks interleave them).
    Otherwise `sort=True` yields geometric reading order (top→bottom, left→right).
    Each text block becomes a paragraph, promoted to a heading when its font is
    materially larger than the page's median span size (or bold + short). Image
    blocks (type 1) are skipped in Phase 1 (figures are the Phase-2/G10 job).
    Returns [] for a text-empty page."""
    import statistics

    from app.document_model import BBox, Block, BlockKind, SOURCE_NATIVE

    try:
        columnar = _page_to_blocks_columnar(page, pno)
        if columnar:
            return columnar
    except Exception as e:  # noqa: BLE001 — column detection must never break ingest
        log.debug("columnar parse failed on page %s (%s); native block parse", pno, e)

    d = page.get_text("dict", sort=True)
    page_w = float(d.get("width") or 0.0)
    page_h = float(d.get("height") or 0.0)

    raw: list[tuple[str, float, bool, tuple]] = []
    sizes: list[float] = []
    for b in d.get("blocks", []):
        if b.get("type") != 0:  # 0 = text, 1 = image (skipped this phase)
            continue
        lines_out: list[str] = []
        max_size = 0.0
        bold = False
        for ln in b.get("lines", []):
            parts = []
            for sp in ln.get("spans", []):
                t = sp.get("text", "")
                parts.append(t)
                max_size = max(max_size, float(sp.get("size", 0.0) or 0.0))
                if int(sp.get("flags", 0) or 0) & _FITZ_BOLD:
                    bold = True
            line = "".join(parts)
            if line.strip():
                lines_out.append(line)
        block_text = "\n".join(lines_out).strip()
        if not block_text:
            continue
        raw.append((block_text, max_size, bold, b.get("bbox")))
        if max_size:
            sizes.append(max_size)

    if not raw:
        return []
    median = statistics.median(sizes) if sizes else 0.0

    out = []
    for text, size, bold, bbox in raw:
        kind, level = BlockKind.PARAGRAPH, 0
        words = len(text.split())
        if median and words <= 20:
            if size >= median * 1.4:
                kind, level = BlockKind.HEADING, 1
            elif size >= median * 1.15 or (bold and words <= 12):
                kind, level = BlockKind.HEADING, 2
        bb = None
        if bbox and page_w and page_h:
            bb = BBox(float(bbox[0]), float(bbox[1]), float(bbox[2]), float(bbox[3]),
                      page_w, page_h, pno)
        out.append(Block(kind=kind, page=pno, text=text, level=level,
                         bbox=bb, source=SOURCE_NATIVE))
    return out


def parse_pdf_structured(pdf_bytes: bytes):
    """Phase 1 · parse a text-layer PDF into the Document Model (IR). Per-page
    errors fall back to a single flat paragraph block; text-empty pages are kept
    (empty paragraph block) so the page index stays aligned — matching parse_pdf.
    Serialise with `.to_pages()` to feed the existing (page, text) contract."""
    from app.document_model import Block, BlockKind, Document, SOURCE_NATIVE

    blocks: list[Block] = []
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for pno, page in enumerate(doc, start=1):
            page_blocks = []
            try:
                page_blocks = _page_to_blocks(page, pno)
            except Exception as e:  # noqa: BLE001 — never let one page kill ingest
                log.warning("doc-model: structured parse failed on page %s (%s); flat fallback", pno, e)
                try:
                    page_blocks = [Block(kind=BlockKind.PARAGRAPH, page=pno,
                                         text=page.get_text("text") or "", source=SOURCE_NATIVE)]
                except Exception:  # noqa: BLE001
                    page_blocks = []
            if not page_blocks:  # preserve the (possibly empty) page slot
                page_blocks = [Block(kind=BlockKind.PARAGRAPH, page=pno, text="", source=SOURCE_NATIVE)]
            blocks.extend(page_blocks)
    return Document(blocks=blocks)


def _parse_pdf_pages(pdf_bytes: bytes) -> list[tuple[int, str]]:
    """PDF → [(page, text)]. Priority: Docling > PyMuPDF structured > flat."""
    s = get_settings()
    if s.doc_model and s.documents_docling_enabled:
        dl = parse_pdf_docling(pdf_bytes)
        if dl is not None:
            return dl[0]
    if s.doc_model:
        try:
            return parse_pdf_structured(pdf_bytes).to_pages()
        except Exception as e:  # noqa: BLE001
            log.warning("doc-model: structured PDF parse failed (%s); using flat parse", e)
    return parse_pdf(pdf_bytes)


def _vision_md_to_pages(pages_md: list[tuple[int, str]]) -> list[tuple[int, str]]:
    """Phase 2 · serialise per-page vision Markdown through the Document Model, so
    key_value / table / heading structure survives (typed blocks → clean text). Falls
    back to the raw Markdown pages if the parse fails — never loses the transcription."""
    try:
        from app import markdown_ir
        return markdown_ir.document_from_pages_markdown(pages_md).to_pages()
    except Exception as e:  # noqa: BLE001
        log.warning("doc-model: markdown→IR failed (%s); using raw vision markdown", e)
        return pages_md


def _table_to_markdown(rows: list[list]) -> str | None:
    """Render a pdfplumber table (list of rows, each a list of cell strings)
    as a GitHub-flavoured Markdown table. Returns None for tables too small
    or too sparse to be worth keeping (pdfplumber over-detects layout columns
    as 'tables', so we filter aggressively)."""
    # Normalise: collapse internal whitespace, treat None as "".
    norm: list[list[str]] = []
    for r in rows or []:
        norm.append([" ".join((str(c) if c is not None else "").split()) for c in r])
    # Drop fully-empty trailing/leading rows.
    norm = [r for r in norm if any(cell for cell in r)]
    if len(norm) < 2:
        return None  # need a header + at least one body row
    width = max(len(r) for r in norm)
    # Drop entirely-empty COLUMNS — the text strategy emits phantom columns for
    # borderless tables ("| | | | CITI | PREMIER |"), which read as noise. Keep
    # only columns that hold a value in at least one row.
    keep = [c for c in range(width) if any(c < len(r) and r[c] for r in norm)]
    if keep and len(keep) < width:
        norm = [[r[c] if c < len(r) else "" for c in keep] for r in norm]
        width = len(keep)
    if width < 2:
        # Single-column "table" → prose paragraph (terms, conditions, notes).
        # Only keep if substantial (≥5 rows) — fewer is layout noise.
        if len(norm) >= 5:
            lines = [c[0] for c in norm if c and c[0]]
            if len(lines) >= 5:
                return "\n".join(lines)
        return None
    # Reject very sparse grids (mostly empty cells = layout noise, not a table).
    filled = sum(1 for r in norm for c in r if c)
    if filled < (len(norm) * width) * 0.4:
        return None
    # Reject LAYOUT tables masquerading as data: a two-column page (e.g. a résumé)
    # gets detected as a grid whose cells each hold a whole section of prose. Real
    # data cells are short (dates, amounts, labels). If any cell is paragraph-sized,
    # or one cell hoards most of the table's text, it's a layout artefact — dropping
    # it stops that scrambled, interleaved text from polluting retrieval.
    cell_lens = [len(c) for r in norm for c in r if c]
    total_chars = sum(cell_lens)
    if cell_lens and (max(cell_lens) > 300 or (total_chars and max(cell_lens) > 0.55 * total_chars)):
        return None
    return _render_md_table(norm, width)


def _render_md_table(rows: list[list[str]], width: int) -> str:
    """Render already-cleaned rows (first row = header) as a GitHub-flavoured
    Markdown table. Rows are padded/trimmed to `width`; pipes inside cells are
    escaped so they don't break the grid. Shared by the pdfplumber table path
    and the CSV parser."""
    def _row(cells: list[str]) -> str:
        padded = (list(cells) + [""] * width)[:width]
        return "| " + " | ".join(c.replace("|", "\\|") for c in padded) + " |"

    header, *body = rows
    lines = [_row(header), "| " + " | ".join(["---"] * width) + " |"]
    lines += [_row(r) for r in body]
    return "\n".join(lines)


# Text-based table strategy — recovers BORDERLESS tables (bank/CC statements,
# financial reports) that the default line-based (lattice) strategy misses
# entirely because there are no ruling lines. Tuned tolerances group columns by
# text alignment. Noisier than lattice, but `_table_to_markdown`'s filters
# (≥2 rows, ≥2 cols, ≥40% filled) drop the layout-fragment false positives.
_TEXT_TABLE_SETTINGS = {
    "vertical_strategy": "text",
    "horizontal_strategy": "text",
    "snap_y_tolerance": 4,
    "intersection_x_tolerance": 12,
}
_MAX_TABLES_PER_PAGE = 10  # raised from 3 — captures more tables per page before density filtering


def _looks_like_header(row: list[str]) -> bool:
    """A row reads as a column header when most non-empty cells are short text
    labels (no leading digit / currency / date) — used to tell a genuine header
    from a continuation page's first DATA row, so we can carry the header over."""
    cells = [c for c in (row or []) if c]
    if len(cells) < 2:
        return False
    labelish = sum(1 for c in cells
                   if not re.match(r"^[\s$€£₹]*[\d.,/-]", c) and len(c) <= 40)
    return labelish >= max(2, (len(cells) + 1) // 2)


def _normalize_table(rows: list[list]) -> list[list[str]]:
    norm = [[" ".join((str(c) if c is not None else "").split()) for c in r] for r in (rows or [])]
    return [r for r in norm if any(cell for cell in r)]


def _camelot_tables(pdf_bytes: bytes) -> list[tuple[list[list[str | None]], int, tuple[float, float, float, float] | None]]:
    """G12 · Try Camelot for table extraction (handles merged cells, multi-page).
    Returns list of (rows, page_number, bbox) tuples where bbox is (x0,y0,x1,y1)
    in PDF points (bottom-left origin). Empty list if Camelot unavailable or no
    tables found. Best-effort — never blocks ingestion."""
    try:
        import camelot
        import tempfile
        import os
        # Camelot needs a file path, not bytes
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
            tf.write(pdf_bytes)
            tmp = tf.name
        try:
            tables = camelot.read_pdf(tmp, pages="all", flavor="lattice")
            if not tables:
                tables = camelot.read_pdf(tmp, pages="all", flavor="stream")
            out = []
            for t in tables:
                if t.data and len(t.data) >= 2:  # header + at least 1 data row
                    bb = None
                    try:
                        bb = tuple(float(v) for v in t._bbox)
                    except Exception:
                        pass
                    out.append((t.data, t.page, bb))
            if out:
                log.info("camelot: extracted %d table(s) across %d page(s)", len(out), len(set(t.page for t in tables)))
            return out
        finally:
            os.unlink(tmp)
    except ImportError:
        log.debug("camelot: not installed — using pdfplumber only")
        return []
    except Exception as e:
        log.debug("camelot: extraction failed, falling back to pdfplumber: %s", e)
        return []


def _split_prose_tail(
    rows: list[list[str | None]],
) -> list[list[list[str | None]]]:
    """Split a table when its tail switches to single-column prose.

    Camelot sometimes merges a data grid (e.g. ticket seating) with an adjacent
    prose block (terms & conditions) because they share column widths.  This
    detects the switch point — rows where most columns are empty and the one
    filled column is long prose text — and splits them into separate tables so
    each PDF box becomes its own chunk + bbox.

    Returns a list of row-groups; the original table when no split is needed."""
    if len(rows) < 2:
        return [rows]

    width = max(len(r) for r in rows)
    if width < 2:
        return [rows]

    # A row is "prose-dominant" when ≥ half the columns are empty AND at least
    # one filled cell is long natural-language text (≥35 chars — shorter than a
    # typical table cell, longer than any label like "Cat 3" or "15%").
    _PROSE_MIN = 35

    split_at: int | None = None
    for i, row in enumerate(rows):
        if i == 0:
            continue  # never split at the header
        filled = [c for c in row if c and str(c).strip()]
        empty = width - len(filled)
        has_prose = any(len(str(c).strip()) >= _PROSE_MIN for c in filled)
        if empty >= width // 2 and has_prose:
            split_at = i
            break

    if split_at is None or split_at < 2:
        return [rows]

    return [rows[:split_at], rows[split_at:]]


def extract_table_chunks(pdf_bytes: bytes) -> list["Chunk"]:
    """P9.5 + G8 + G12 · Detect tables via Camelot → pdfplumber → Markdown chunks
    (kind='table'). Camelot handles merged cells and multi-page tables; pdfplumber
    is the proven fallback for borderless/irregular layouts.

    After extraction, adjacent-but-distinct blocks that Camelot merged are split
    back apart: rows that switch from structured data to single-column prose
    (terms, conditions, notes) are separated into their own chunk so they don't
    visually merge two PDF boxes into one highlight band.

    Best-effort: any failure logs and returns whatever was gathered — never
    blocks ingestion.
    """
    import pdfplumber

    out: list[Chunk] = []
    last_header: dict[int, tuple[list[str], int]] = {}

    # G12 · Try Camelot first (handles merged cells, multi-page tables natively)
    camelot_raw = _camelot_tables(pdf_bytes)
    if camelot_raw:
        for data, page_no, camelot_bb in camelot_raw:
            for rows in _split_prose_tail(data):
                md = _table_to_markdown(rows)
                if not md:
                    continue
                is_table = md.startswith("|")
                label = f"[{'Table' if is_table else 'Text'} · page {page_no}]"
                text = f"{label}\n{md}"
                out.append(Chunk(
                    text=text, page=page_no, char_start=0, char_end=len(text),
                    kind="table",  # from extractor — uses page-based bbox fallback
                    table_bbox=camelot_bb,  # for spatial line matching
                ))
        if out:
            log.info("tables: %d chunk(s) from Camelot", len(out))
            return out

    # Fallback: pdfplumber (lattice + text strategy)
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page_no, page in enumerate(pdf.pages, start=1):
                try:
                    raw = page.extract_tables() or []
                except Exception as e:
                    log.debug("pdfplumber: page %d lattice failed: %s", page_no, e)
                    raw = []
                tables = [t for t in (_normalize_table(t) for t in raw) if len(t) >= 2]
                if not tables:
                    try:
                        raw = page.extract_tables(_TEXT_TABLE_SETTINGS) or []
                    except Exception as e:
                        log.debug("pdfplumber: page %d text-strategy failed: %s", page_no, e)
                        raw = []
                    tables = [t for t in (_normalize_table(t) for t in raw) if len(t) >= 2]
                    tables.sort(key=lambda t: sum(1 for r in t for c in r if c), reverse=True)
                    tables = tables[:_MAX_TABLES_PER_PAGE]

                for t in tables:
                    t_norm = _normalize_table(t)
                    for rows in _split_prose_tail(t_norm):
                        width = max(len(r) for r in rows)
                        cont = False
                        prior = last_header.get(width)
                        if prior and prior[1] == page_no - 1 and not _looks_like_header(rows[0]):
                            rows = [prior[0]] + rows
                            last_header[width] = (prior[0], page_no)
                            cont = True
                        elif _looks_like_header(rows[0]):
                            last_header[width] = (rows[0], page_no)
                        md = _table_to_markdown(rows)
                        if not md:
                            continue
                        label = f"[Table · page {page_no}{' (cont.)' if cont else ''}]\n"
                        text = label + md
                        out.append(Chunk(
                            text=text, page=page_no,
                            char_start=0, char_end=len(text), kind="table",
                        ))
    except Exception as e:
        log.warning("pdfplumber table extraction failed (non-fatal): %s", e)
    if out:
        log.info("pdfplumber: extracted %d table chunk(s)", len(out))
    return out


class _ImgOcrBudget:
    """Shared state for embedded-image OCR across one document: the opt-in flag,
    a per-document call budget, and a dedup set. For a multi-slide deck the same
    budget threads through every slide so the cap is per-document, not per-slide."""

    def __init__(self):
        settings = get_settings()
        self.enabled: bool = is_enabled("documents_office_image_ocr", False)
        self.remaining: int = max(0, get_int("documents_office_image_max", 12))
        self.seen: set[str] = set()


def _ocr_office_images(
    blobs, budget: _ImgOcrBudget, *, db: Session | None, tenant_id: str | None,
) -> list[str]:
    """Vision-OCR a stream of `(blob, content_type)` embedded images, drawing down
    the shared per-document `budget`. Returns one `[Image OCR] …` line per image
    that yielded text. De-dups repeated images (a screenshot reused across slides
    is read once) and skips decorative chrome (logos/icons) before spending a
    vision call. Disabled or exhausted → []."""
    if not budget.enabled:
        return []
    out: list[str] = []
    for blob, ctype in blobs:
        if budget.remaining <= 0:
            break
        if not blob:
            continue
        h = hashlib.md5(blob).hexdigest()  # noqa: S324 — dedup key, not security
        if h in budget.seen:
            continue
        budget.seen.add(h)
        if not ingestion_vision.embedded_image_is_content(blob):
            continue  # logo/icon/bullet — no VLM call, no budget spent
        txt = ingestion_vision.ocr_embedded_image(
            blob, ctype, db=db, tenant_id=tenant_id,
        )
        budget.remaining -= 1  # a vision call was made, text or not
        if txt and txt.strip():
            out.append(f"[Image OCR] {txt.strip()}")
    return out


def _docx_image_blobs(d):
    """Yield `(blob, content_type)` for each image embedded in the docx body.
    Skips external (linked) images. Header/footer parts aren't walked, which
    conveniently keeps letterhead logos out of the OCR budget."""
    for rel in list(d.part.rels.values()):
        try:
            if rel.is_external or "image" not in rel.reltype:
                continue
            part = rel.target_part
            yield part.blob, getattr(part, "content_type", None)
        except Exception:  # noqa: BLE001 — tolerate a broken relationship
            continue


def parse_docx(
    raw: bytes, *, db: Session | None = None, tenant_id: str | None = None,
    return_ir: bool = False,
) -> list[tuple[int, str]] | tuple[list[tuple[int, str]], object]:
    """T1.2 · Parse a .docx via python-docx. Returns pages, or (pages, ir_document)
    when return_ir=True and doc_model is enabled. IR preserves heading hierarchy,
    table structure, and list items for block-aware chunking."""
    # Guard against decompression-bomb DoS (authenticated per-user, but 100 MB zip →
    # GB RAM via lxml full-load). 50 MB cap is generous for real office docs.
    if len(raw) > 50 * 1024 * 1024:
        raise ValueError(f"DOCX too large ({len(raw) / 1024 / 1024:.1f} MB); max 50 MB")
    from io import BytesIO
    from docx import Document as DocxDocument

    out_parts: list[str] = []
    d = DocxDocument(BytesIO(raw))
    for para in d.paragraphs:
        t = (para.text or "").strip()
        if t:
            out_parts.append(t)
    for table in d.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            cells = [c for c in cells if c]
            if cells:
                out_parts.append(" | ".join(cells))
    img_lines = _ocr_office_images(
        _docx_image_blobs(d), _ImgOcrBudget(), db=db, tenant_id=tenant_id,
    )
    if get_settings().doc_model:
        ir = _docx_to_document(d, img_lines)
        if return_ir:
            return ir.to_pages(), ir
        return ir.to_pages()
    out_parts.extend(img_lines)
    text = "\n\n".join(out_parts).strip()
    pages = [(1, text)] if text else []
    if return_ir:
        return pages, None
    return pages


def _docx_to_document(d, img_lines: list[str]):
    """Phase 3 · map a python-docx document to the IR in body order: heading styles
    → heading blocks (level from 'Heading N'), list styles → list_item, w:tbl → table
    blocks (real cells + header), other paragraphs → paragraph; OCR'd embedded images
    → figure blocks. Everything on 'page 1' (Word has no fixed page concept)."""
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table as _Table
    from docx.text.paragraph import Paragraph as _Para

    from app.document_model import Block, BlockKind, Document, SOURCE_NATIVE, SOURCE_VISION

    blocks: list[Block] = []
    for child in d.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = _Para(child, d)
            text = (para.text or "").strip()
            if not text:
                continue
            style = ((para.style.name if para.style else "") or "")
            if style.startswith("Heading") or style == "Title":
                m = re.search(r"\d+", style)
                blocks.append(Block(kind=BlockKind.HEADING, page=1, text=text,
                                    level=int(m.group()) if m else 1, source=SOURCE_NATIVE))
            elif "List" in style:
                blocks.append(Block(kind=BlockKind.LIST_ITEM, page=1, text=text, source=SOURCE_NATIVE))
            else:
                blocks.append(Block(kind=BlockKind.PARAGRAPH, page=1, text=text, source=SOURCE_NATIVE))
        elif isinstance(child, CT_Tbl):
            tbl = _Table(child, d)
            rows = [[(c.text or "").strip() for c in row.cells] for row in tbl.rows]
            rows = [r for r in rows if any(r)]
            if rows:
                blocks.append(Block(kind=BlockKind.TABLE, page=1, rows=rows,
                                    has_header=True, source=SOURCE_NATIVE))
    for line in img_lines:  # "[Image OCR] …" → figure block
        blocks.append(Block(kind=BlockKind.FIGURE, page=1,
                            ocr_text=line.replace("[Image OCR] ", "", 1), source=SOURCE_VISION))
    return Document(blocks=blocks)


def parse_xlsx(raw: bytes) -> list[tuple[int, str]]:
    """T1.2 · Parse an .xlsx via openpyxl. Each sheet becomes a 'page'.
    Per row: tab-joined cell values. Skips empty rows. Truncates per-cell
    values to 500 chars to keep retrieval chunks sane.
    """
    from io import BytesIO
    from openpyxl import load_workbook

    def _cell(v) -> str:
        if v is None:
            return ""
        s = str(v).strip()
        return (s[:500] + "…") if len(s) > 500 else s

    wb = load_workbook(BytesIO(raw), read_only=True, data_only=True)
    structured = get_settings().doc_model
    if structured:
        from app.document_model import Block, BlockKind, Document, SOURCE_NATIVE
        blocks: list[Block] = []
        for sheet_idx, ws in enumerate(wb.worksheets, start=1):
            name = ws.title or f"Sheet {sheet_idx}"
            rows = [[_cell(v) for v in row] for row in ws.iter_rows(values_only=True)]
            rows = [r for r in rows if any(r)]
            if not rows:
                continue
            blocks.append(Block(kind=BlockKind.HEADING, page=sheet_idx,
                                text=f"Sheet: {name}", level=1, source=SOURCE_NATIVE))
            blocks.append(Block(kind=BlockKind.TABLE, page=sheet_idx, rows=rows,
                                has_header=True, source=SOURCE_NATIVE))
        return Document(blocks=blocks).to_pages()

    out: list[tuple[int, str]] = []
    for sheet_idx, ws in enumerate(wb.worksheets, start=1):
        lines: list[str] = [f"# Sheet: {ws.title or f'Sheet {sheet_idx}'}"]
        for row in ws.iter_rows(values_only=True):
            cells = [_cell(v) for v in row]
            if any(c for c in cells):
                lines.append("\t".join(cells))
        text = "\n".join(lines).strip()
        if text:
            out.append((sheet_idx, text))
    return out


def parse_text(raw: bytes) -> list[tuple[int, str]]:
    """T1.2 · Plain text or markdown. Uses the same robust decoder as CSV
    (chardet → utf-8-sig/utf-8/cp1252/latin-1/utf-16) so a Windows "Unicode"
    (UTF-16LE) .txt decodes cleanly instead of garbling to interleaved-NUL
    Latin-1. `normalize_text` downstream strips any residual BOM/control bytes."""
    text = _decode_csv(raw).strip()
    return [(1, text)] if text else []


def parse_csv(raw: bytes) -> list[tuple[int, str]]:
    """Parse a CSV/TSV into a Markdown table (first row = header) so retrieval and
    the fact extractor see real tabular structure — not raw delimited bytes.

    Uses the stdlib csv reader, so it correctly handles what a raw-text passthrough
    garbled: quoted fields, embedded delimiters/newlines (a quoted cell with a
    newline no longer splits into two rows), and non-comma delimiters (`;` European
    exports, `\\t` TSV, `|`) via delimiter sniffing. Cell values are collapsed to a
    single line and capped at 500 chars (matching parse_xlsx). Falls back to the
    plain decoded text if the content can't be read as a table."""
    import csv as _csv

    text = _decode_csv(raw)
    if not text.strip():
        return []
    # Sniff the delimiter from a sample; restrict candidates to reduce misdetection.
    try:
        delimiter = _csv.Sniffer().sniff(text[:8192], delimiters=",;\t|").delimiter
    except Exception:  # noqa: BLE001 — unsniffable (single column / irregular)
        delimiter = ","
    try:
        parsed = list(_csv.reader(io.StringIO(text), delimiter=delimiter))
    except Exception:  # noqa: BLE001 — malformed; keep the plain text
        return [(1, text.strip())]
    rows: list[list[str]] = []
    for r in parsed:
        # Collapse whitespace (incl. embedded newlines) so each row is one line;
        # cap cell width like parse_xlsx to keep retrieval chunks sane.
        cells = [" ".join((c or "").split()) for c in r]
        cells = [(c[:500] + "…") if len(c) > 500 else c for c in cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return []
    width = max(len(r) for r in rows)
    md = _render_md_table(rows, width) if width >= 1 else ""
    return [(1, md)] if md.strip() else [(1, text.strip())]


def _html_to_text(html: str) -> str:
    """Reduce an HTML email body to readable text: drop <script>/<style> blocks
    ENTIRELY (their JS/CSS contents are not body text), strip remaining tags, then
    decode HTML entities (&amp; &nbsp; &#39; …). A regex reducer, not a renderer —
    good enough for indexing; avoids the wall-of-CSS-noise the old tag-only strip
    produced for style-heavy marketing/newsletter emails."""
    import html as _htmlmod
    import re
    s = re.sub(r"(?is)<(script|style)\b[^>]*>.*?</\1>", " ", html or "")
    s = re.sub(r"(?s)<[^>]+>", " ", s)
    s = _htmlmod.unescape(s)
    return re.sub(r"\s+", " ", s).strip()


def parse_eml(raw: bytes) -> list[tuple[int, str]]:
    """T1.2 · RFC 5322 email. Emits structured 'page 1' with headers
    (From / To / Subject / Date) followed by the text/plain body. HTML
    parts get HTML tags stripped via a simple regex — good enough for
    body indexing, not a full HTML renderer.
    """
    import email
    from email import policy

    msg = email.message_from_bytes(raw, policy=policy.default)
    headers = []
    for h in ("From", "To", "Cc", "Subject", "Date"):
        v = msg.get(h)
        if v:
            headers.append(f"{h}: {v}")

    # Pick best body part.
    body = ""
    if msg.is_multipart():
        # Prefer text/plain over text/html.
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                try:
                    body = part.get_content()
                    break
                except Exception:
                    pass
        if not body:
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    try:
                        body = _html_to_text(part.get_content())
                        break
                    except Exception:
                        pass
    else:
        try:
            body = msg.get_content() or ""
        except Exception:
            body = msg.get_payload(decode=True)
            if isinstance(body, bytes):
                body = body.decode("utf-8", errors="replace")
            body = body or ""
        # Single-part text/html email → clean it the same way as the multipart branch.
        if msg.get_content_type() == "text/html" and body:
            body = _html_to_text(body)

    text = "\n".join(headers) + "\n\n" + (body or "")
    text = text.strip()
    return [(1, text)] if text else []


def _pptx_picture_blobs(shapes):
    """Yield `(blob, content_type)` for every picture on a slide, descending into
    group shapes. Reads `shape.image` inside a guard so it captures both plain
    pictures and picture placeholders while ignoring text/table/vector shapes
    (which raise on `.image`)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:  # noqa: BLE001 — python-pptx layout changed; degrade to none
        return
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _pptx_picture_blobs(shape.shapes)
                continue
        except Exception:  # noqa: BLE001
            pass
        try:
            img = shape.image  # Picture / picture placeholder; raises otherwise
            yield img.blob, img.content_type
        except Exception:  # noqa: BLE001 — not a picture, skip
            continue


def parse_pptx(
    raw: bytes, *, db: Session | None = None, tenant_id: str | None = None,
    return_ir: bool = False,
) -> list[tuple[int, str]] | tuple[list[tuple[int, str]], object]:
    """Parse a .pptx via python-pptx. Each slide becomes a 'page'. When
    return_ir=True and doc_model is enabled, returns (pages, ir_document)
    for block-aware chunking."""
    # Guard against decompression-bomb DoS (same rationale as parse_docx).
    if len(raw) > 50 * 1024 * 1024:
        raise ValueError(f"PPTX too large ({len(raw) / 1024 / 1024:.1f} MB); max 50 MB")
    from io import BytesIO

    from pptx import Presentation  # python-pptx

    prs = Presentation(BytesIO(raw))
    structured = get_settings().doc_model
    if structured:
        from app.document_model import Block, BlockKind, Document, SOURCE_NATIVE, SOURCE_VISION
    pages: list[tuple[int, str]] = []
    blocks: list = []
    img_budget = _ImgOcrBudget()  # one per-document budget across all slides
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:  # noqa: BLE001
            pass
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    if structured:
                        if shape is title_shape:
                            blocks.append(Block(kind=BlockKind.HEADING, page=i, text=t, level=2, source=SOURCE_NATIVE))
                        else:
                            blocks.append(Block(kind=BlockKind.PARAGRAPH, page=i, text=t, source=SOURCE_NATIVE))
                    else:
                        parts.append(t)
            if shape.has_table:
                rows = [[(c.text or "").strip() for c in row.cells] for row in shape.table.rows]
                rows = [r for r in rows if any(r)]
                if structured:
                    if rows:
                        blocks.append(Block(kind=BlockKind.TABLE, page=i, rows=rows, has_header=True, source=SOURCE_NATIVE))
                else:
                    for r in rows:
                        parts.append(" | ".join(c for c in r if c))
        # Speaker notes — often carry the substance of a deck.
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    if structured:
                        blocks.append(Block(kind=BlockKind.PARAGRAPH, page=i, text=f"[Notes] {note}", source=SOURCE_NATIVE))
                    else:
                        parts.append(f"[Notes] {note}")
        except Exception:  # noqa: BLE001
            pass
        img_lines = _ocr_office_images(
            _pptx_picture_blobs(slide.shapes), img_budget, db=db, tenant_id=tenant_id,
        )
        if structured:
            for line in img_lines:
                blocks.append(Block(kind=BlockKind.FIGURE, page=i,
                                    ocr_text=line.replace("[Image OCR] ", "", 1), source=SOURCE_VISION))
        else:
            parts.extend(img_lines)
            text = "\n\n".join(parts).strip()
            if text:
                pages.append((i, text))
    if structured:
        ir = Document(blocks=blocks)
        if return_ir:
            return ir.to_pages(), ir
        return ir.to_pages()
    if return_ir:
        return pages, None
    return pages


# Formats LibreOffice converts → PDF (then parse_pdf handles them). Opt-in;
# only fires when soffice is installed (full image), no-op otherwise.
_OFFICE_CONVERT_EXTS = {
    ".doc", ".ppt", ".xls",            # legacy MS Office (OLE)
    ".odt", ".odp", ".ods",            # OpenDocument
    ".rtf",                            # Rich Text
}


def libreoffice_to_pdf(raw: bytes, ext: str) -> bytes | None:
    """Convert an office doc to PDF via headless LibreOffice, HARDENED. Returns
    PDF bytes or None (caller falls back / fails). No-op + None when soffice is
    absent (slim image) — same optional pattern as the rapidocr engine.

    Hardening (untrusted-doc conversion is a real attack surface):
      · macros never execute (--convert-to doesn't, and we don't enable them)
      · a throwaway, isolated user profile per call (-env:UserInstallation)
      · NO network (LibreOffice headless convert makes none; the container has
        no egress need for this path)
      · hard wall-clock timeout so a malicious file can't hang the worker
      · runs as the existing non-root container user
    """
    import shutil
    import subprocess
    import tempfile

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None  # not installed (e.g. slim image) — optional path
    try:
        with tempfile.TemporaryDirectory() as td:
            src = os.path.join(td, f"in{ext}")
            with open(src, "wb") as f:
                f.write(raw)
            profile = os.path.join(td, "profile")
            cmd = [
                soffice, "--headless", "--norestore", "--nolockcheck",
                f"-env:UserInstallation=file://{profile}",
                "--convert-to", "pdf", "--outdir", td, src,
            ]
            subprocess.run(cmd, cwd=td, timeout=120, check=True,
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            out_pdf = os.path.join(td, "in.pdf")
            if os.path.exists(out_pdf):
                with open(out_pdf, "rb") as f:
                    return f.read()
    except Exception as e:  # noqa: BLE001 — never block ingest on a conversion failure
        log.warning("libreoffice convert failed for %s (non-fatal): %s", ext, e)
    return None


_MRZ_LINE_RE = re.compile(r"[A-Z0-9<]{30,44}")


def _mrz_span(text: str) -> tuple[int, int] | None:
    """T2.3 · Locate a 2-line MRZ block in text (passport MRZ Type 1/3 = 44
    chars × 2 lines; ID card MRZ Type 1 = 30 chars × 3 lines). Returns
    (start, end) char range covering both/all lines, or None when no MRZ
    is present. Used by the chunker to avoid splitting an MRZ across two
    chunks — which would break grounding for passport_no / surname /
    given_name / DOB / nationality."""
    matches = list(_MRZ_LINE_RE.finditer(text))
    if len(matches) < 2:
        return None
    # Adjacent or near-adjacent matches (within 5 chars of whitespace
    # between them) — heuristic for "same MRZ block".
    for i in range(len(matches) - 1):
        a = matches[i]
        b = matches[i + 1]
        gap = b.start() - a.end()
        if gap < 10 and len(a.group()) >= 30 and len(b.group()) >= 30:
            # Look ahead for an optional 3rd line (ID card format).
            end = b.end()
            if i + 2 < len(matches):
                c = matches[i + 2]
                if c.start() - b.end() < 10 and len(c.group()) >= 30:
                    end = c.end()
            return (a.start(), end)
    return None


def chunk_pages(pages: list[tuple[int, str]]) -> list[Chunk]:
    """Layout-aware chunker (G5). Each chunk carries the page it started on.

    Delegates the per-page algorithm to `app.chunking.chunk_page_text`, which
    splits the page into paragraph blocks and packs WHOLE blocks into
    ~target-sized chunks instead of cutting fixed char windows mid-sentence.

    T2.3 · MRZ-aware: a passport/ID MRZ block stays intact — small pages become
    one chunk (MRZ whole), an MRZ paragraph stays a whole block, and if an
    oversized block must be windowed the MRZ span is passed as a protected span
    so the surname/given-name parsing + tight-bbox grounding keep working.
    """
    from app import chunking

    settings = get_settings()
    target = settings.chunk_target_chars
    overlap = settings.chunk_overlap_chars
    out: list[Chunk] = []

    # Per-page chunk cap. The old default (50) silently dropped a page's tail once
    # exceeded — invisible content loss for large SINGLE-page sources (a big .txt /
    # .csv / .eml is emitted as one "page"). Raise it generously and LOG when a page
    # actually hits the cap so truncation is never silent.
    MAX_PER_PAGE = 400
    for page_no, page_text in pages:
        # R6 · NFKC normalize before chunking (safe, deterministic).
        if settings.chunk_nfkc_normalize:
            page_text = chunking.normalize_text(page_text)
        page_chunks = chunking.chunk_page_text(
            page_text, target, overlap, protect_span_fn=_mrz_span,
            max_chunks=MAX_PER_PAGE,
            sentence_aware=settings.semantic_chunking,   # RAG-roadmap #4 (default off)
        )
        if len(page_chunks) >= MAX_PER_PAGE:
            log.warning("chunk_pages: page %s hit the %d-chunk cap — tail may be truncated (len=%d chars)",
                        page_no, MAX_PER_PAGE, len(page_text or ""))
        for text, char_start, char_end in page_chunks:
            out.append(
                Chunk(text=text, page=page_no, char_start=char_start, char_end=char_end)
            )

    # R6 · drop near-duplicate chunks (recurring boilerplate across pages).
    if settings.chunk_dedup_near_duplicates and len(out) > 1:
        keep = set(chunking.dedup_indices([c.text for c in out], settings.chunk_dedup_threshold))
        if len(keep) < len(out):
            log.info("chunk_pages: deduped %d near-duplicate chunk(s)", len(out) - len(keep))
            out = [c for i, c in enumerate(out) if i in keep]
    return out


def chunk_document(document) -> list[Chunk]:
    """Phase 4 · block-aware chunking straight from the Document Model (IR). Groups
    blocks by page and delegates to `chunking.chunk_blocks`, which never splits a
    key_value/heading block and row-splits an oversized table with the header
    repeated. Sanitises + NFKC-normalises each block's text (the flat path does this
    per-page); then the same near-duplicate dedup as chunk_pages."""
    from collections import OrderedDict

    from app import chunking

    settings = get_settings()
    target, overlap = settings.chunk_target_chars, settings.chunk_overlap_chars
    nfkc = settings.chunk_nfkc_normalize
    MAX_PER_PAGE = 400

    def _clean(s: str) -> str:
        s = _sanitize_text(s)
        return chunking.normalize_text(s) if nfkc else s

    # Tag every block with its document-level index so chunk_blocks can recover
    # the mapping from chunk → block IDs (b_0000, b_0001, …) for block_map.
    for i, b in enumerate(document.blocks):
        b._doc_idx = i

    by_page: "OrderedDict[int, list]" = OrderedDict()
    for b in document.blocks:
        by_page.setdefault(b.page, []).append(b)

    out: list[Chunk] = []
    for page_no, blocks in by_page.items():
        for b in blocks:                       # clean text fields in place (transient IR)
            if getattr(b, "text", ""):
                b.text = _clean(b.text)
            if getattr(b, "label", ""):
                b.label = _clean(b.label)
            if getattr(b, "value", ""):
                b.value = _clean(b.value)
            if getattr(b, "rows", None):
                b.rows = [[_clean(str(c)) for c in r] for r in b.rows]
            for attr in ("caption", "description", "ocr_text"):
                if getattr(b, attr, ""):
                    setattr(b, attr, _clean(getattr(b, attr)))
        page_chunks = chunking.chunk_blocks(blocks, target, overlap, max_chunks=MAX_PER_PAGE)
        if len(page_chunks) >= MAX_PER_PAGE:
            log.warning("chunk_document: page %s hit the %d-chunk cap", page_no, MAX_PER_PAGE)
        for text, cs, ce, bbox, block_indices in page_chunks:
            # Convert document-level integer indices to string block IDs
            bid_tuple = tuple(f"b_{bi:04d}" for bi in block_indices) if block_indices else ()
            out.append(Chunk(text=text, page=page_no, char_start=cs, char_end=ce,
                             bbox=bbox, block_ids=bid_tuple))

    if settings.chunk_dedup_near_duplicates and len(out) > 1:
        keep = set(chunking.dedup_indices([c.text for c in out], settings.chunk_dedup_threshold))
        if len(keep) < len(out):
            log.info("chunk_document: deduped %d near-duplicate chunk(s)", len(out) - len(keep))
            out = [c for i, c in enumerate(out) if i in keep]
    return out


def _chunk_annotated_markdown(
    md: str, target: int, overlap: int, *, nfkc: bool = True,
) -> list[Chunk]:
    """Chunk annotated markdown (with ``<!-- block:b_XXXX -->`` markers) into
    Chunk objects that carry their composing ``block_ids``.

    This is the reprocess-path counterpart of ``chunk_document``: instead of
    the Docling IR it parses the block markers the frontend preserves, joins
    segment texts the same way (``\\n\\n``), runs the same ``chunk_page_text``
    packer, then maps each chunk's char span back to the segment(s) it overlaps.

    Returns chunks with ``page=1`` (markdown doesn't carry page info) and
    ``block_ids`` populated (empty tuple for segments without a marker)."""
    import re as _re

    from app import chunking as _chunking

    # 1. Split on block markers → interleaved [text, marker, text, marker, …]
    BLOCK_RX = _re.compile(r"<!-- block:(b_\w+) -->")
    parts = BLOCK_RX.split(md or "")
    segments: list[tuple[str | None, str]] = []  # (block_id | None, text)
    i = 0
    while i < len(parts):
        part = parts[i]
        # Every even index is text; odd indices are captured block IDs
        if i % 2 == 0:
            txt = part.strip()
            if txt:
                segments.append((None, txt))
        else:
            # part is the captured block ID
            bid = part
            # next part is the text that follows this marker
            i += 1
            txt = parts[i].strip() if i < len(parts) else ""
            if txt:
                segments.append((bid, txt))
        i += 1

    if not segments:
        return []

    # 2. Normalise each segment text (same pipeline as chunk_document → _clean)
    for j, (bid, txt) in enumerate(segments):
        s = _sanitize_text(txt)
        s = _chunking.normalize_text(s) if nfkc else s
        segments[j] = (bid, s)

    # 3. Join with "\n\n" (same join chunk_blocks uses) and track offsets
    seg_starts: list[int] = []
    pos = 0
    pieces: list[str] = []
    for j, (bid, txt) in enumerate(segments):
        seg_starts.append(pos)
        pieces.append(txt)
        pos += len(txt) + (2 if j < len(segments) - 1 else 0)
    joined = "\n\n".join(pieces)

    # 4. Chunk the joined page string with the same packer
    spans = _chunking.chunk_page_text(joined, target, overlap)
    if not spans:
        return []

    # 5. Map each chunk span back to the segments it overlaps
    out: list[Chunk] = []
    for text, cs, ce in spans:
        chunk_bids: list[str] = []
        for j, (bid, txt) in enumerate(segments):
            seg_start = seg_starts[j]
            seg_end = seg_start + len(txt)
            # Overlap check: chunk [cs, ce) intersects segment [seg_start, seg_end)
            if cs < seg_end and ce > seg_start and bid is not None:
                chunk_bids.append(bid)
        out.append(Chunk(
            text=text, page=1, char_start=cs, char_end=ce,
            block_ids=tuple(chunk_bids) if chunk_bids else (),
        ))
    return out


def embed_chunks(chunks: list[Chunk]) -> list[list[float]]:
    """Embed in modest batches so OpenAI's per-call payload stays reasonable.
    The hash backend doesn't care about batching but uses the same path."""
    BATCH = 64
    out: list[list[float]] = []
    for i in range(0, len(chunks), BATCH):
        out.extend(embed([c.text for c in chunks[i : i + BATCH]]))
    return out


def _extract_text_layer(pdf_bytes: bytes, page_count: int) -> list[dict]:
    """M47 · Extract word-level text blocks from every page for reverse bbox lookup.
    Returns list of {page, x0, y0, x1, y1, page_w, page_h, text, kind}.
    Keeps words, not paragraphs — fine-grained click targeting. Capped at 200 pages."""
    import fitz
    blocks = []
    max_pages = min(page_count, 200)
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for pg in range(max_pages):
                page = doc[pg]
                pw, ph = page.rect.width, page.rect.height
                # Get word-level dict
                words = page.get_text("words")
                for w in words:
                    if len(w) >= 5 and w[4].strip():  # x0,y0,x1,y1,text,block,line,word
                        blocks.append({
                            "page": pg + 1,
                            "x0": round(w[0], 1), "y0": round(w[1], 1),
                            "x1": round(w[2], 1), "y1": round(w[3], 1),
                            "page_w": round(pw, 1), "page_h": round(ph, 1),
                            "text": w[4].strip(),
                            "kind": "word",
                        })
    except Exception:
        pass
    return blocks


def _bboxes_for_chunks(pdf_bytes: bytes, chunks: list[Chunk]) -> list[dict | None]:
    """For each chunk, find a bounding rectangle on its page that covers the
    FULL chunk text span (not just the first word). Uses PyMuPDF word-level
    text extraction to locate all words in the chunk and unions their bboxes.
    Falls back to search_for(text[:80]) when the word-level path can't match."""
    bboxes: list[dict | None] = []
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            page_cache: dict[int, fitz.Page] = {}
            for c in chunks:
                page_idx = c.page - 1
                if page_idx not in page_cache:
                    try:
                        page_cache[page_idx] = doc.load_page(page_idx)
                    except Exception:  # noqa: BLE001
                        bboxes.append(None)
                        continue
                page = page_cache[page_idx]
                ct = (c.text or "").strip()
                if not ct:
                    bboxes.append(None)
                    continue
                pw = float(page.rect.width)
                ph = float(page.rect.height)
                bbox = _locate_text_span(page, ct, pw, ph)
                if bbox:
                    bbox["page"] = c.page
                bboxes.append(bbox)
    except Exception as e:  # noqa: BLE001
        log.warning("bbox extraction failed: %s · falling back to NULLs", e)
        return [None] * len(chunks)
    return bboxes


# ── Line-ID bbox pipeline ──────────────────────────────────────────────────
# Instead of post-hoc word-matching (_locate_text_span), we capture per-line
# geometry at parse time via page.get_text("dict") and carry line_ids through
# chunks. Chunk bbox = union of full-width line bands — no x-coordinate
# matching across columns.

def _build_line_map(pdf_bytes: bytes) -> dict[str, dict]:
    """Build a document-global line_map from a native PDF.

    Returns {line_id: {page, y0_pct, h_pct, page_w, page_h, text}} where
    line_id is a hex string "0x<pno>:<line_ix>" (e.g. "0x1:3f" = page 1,
    line 63).  ``text`` is the joined span text so downstream consumers
    (``_compute_chunk_line_ids``) don't need to re-open the PDF.
    Full-width line bands — x spans the entire page width, so only y matters.
    Percentages of native page dimensions (zoom/resolution independent)."""
    line_map: dict[str, dict] = {}
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for pno, page in enumerate(doc, start=1):
                pw = float(page.rect.width)
                ph = float(page.rect.height)
                if pw <= 0 or ph <= 0:
                    continue
                d = page.get_text("dict", sort=True)
                line_ix = 0
                for b in d.get("blocks", []):
                    if b.get("type") != 0:  # text blocks only
                        continue
                    for ln in b.get("lines", []):
                        bbox = ln.get("bbox")
                        if not bbox or len(bbox) < 4:
                            continue
                        y0, y1 = float(bbox[1]), float(bbox[3])
                        if y1 <= y0:
                            continue
                        lid = f"0x{pno:x}:{line_ix:x}"
                        parts = [sp.get("text", "") for sp in ln.get("spans", [])]
                        line_text = " ".join(parts).strip()
                        line_map[lid] = {
                            "page": pno,
                            "y0_pct": round(y0 / ph, 6),
                            "h_pct": round((y1 - y0) / ph, 6),
                            "page_w": round(pw, 1),
                            "page_h": round(ph, 1),
                            "text": line_text,
                        }
                        line_ix += 1
    except Exception as e:  # noqa: BLE001
        log.warning("line_map: build failed: %s", e)
        return {}
    return line_map


def _match_lines_spatial(
    chunk: Chunk,
    page_lines: list[tuple[str, dict]],
) -> list[str] | None:
    """Match a chunk to line_ids by spatial overlap with its Camelot bbox.

    Camelot bbox is (x0, y0, x1, y1) in PDF points, bottom-left origin.
    PyMuPDF line_map uses top-left origin percentages.  We convert by flipping
    the y-axis:  y_top = page_h - y1,  y_bot = page_h - y0."""
    bb = chunk.table_bbox
    if not bb or not page_lines:
        return None

    # Get page dimensions from any line on this page
    page_h = page_lines[0][1].get("page_h", 0)
    if page_h <= 0:
        return None

    # Camelot bbox in PDF points (bottom-left origin)
    _, y0_cam, _, y1_cam = bb  # y0 < y1 (bottom < top in PDF coords)

    # Convert to top-left percentages (PyMuPDF convention)
    table_y0_pct = (page_h - y1_cam) / page_h  # top of table
    table_y1_pct = (page_h - y0_cam) / page_h  # bottom of table

    if table_y0_pct >= table_y1_pct or table_y0_pct < 0:
        return None

    # Small margin to catch adjacent lines (±2% of page height)
    margin = 0.02
    y0 = max(0.0, table_y0_pct - margin)
    y1 = min(1.0, table_y1_pct + margin)

    lids = []
    for lid, info in page_lines:
        ly0 = info["y0_pct"]
        ly1 = ly0 + info["h_pct"]
        # Line overlaps the table's y-range
        if ly1 >= y0 and ly0 <= y1:
            lids.append(lid)
    return lids if lids else None


def _bbox_from_table_bbox(
    chunk: Chunk, line_map: dict[str, dict]
) -> dict | None:
    """Compute a chunk bbox directly from the Camelot table_bbox, which has
    accurate x-coordinates the full-width line bands don't carry.

    Returns {page, x0_pct, x1_pct, y0_pct, y1_pct, page_w, page_h} or None.
    The y-range still comes from line_map (per-line precision); the x-range
    comes from Camelot's column-aware table detection, clamped sensibly."""
    bb = chunk.table_bbox
    if not bb or not line_map:
        return None
    # Get page dimensions from a line_map entry on the chunk's OWN page
    cpage = chunk.page or 1
    page_w = page_h = 0
    for _info in line_map.values():
        if _info.get("page") == cpage:
            page_w = _info.get("page_w", 0)
            page_h = _info.get("page_h", 0)
            break
    if page_w <= 0 or page_h <= 0:
        # Fallback: any page (single-page docs)
        any_line = next(iter(line_map.values()))
        page_w = any_line.get("page_w", 0)
        page_h = any_line.get("page_h", 0)
    if page_w <= 0 or page_h <= 0:
        return None

    # Camelot bbox in PDF points, bottom-left origin
    x0_cam, y0_cam, x1_cam, y1_cam = bb  # y0<y1 (bottom<top in PDF)

    # Convert to top-left percentages
    x0_pct = max(0.0, min(1.0, x0_cam / page_w))
    x1_pct = max(0.0, min(1.0, x1_cam / page_w))
    table_y0_pct = (page_h - y1_cam) / page_h  # top
    table_y1_pct = (page_h - y0_cam) / page_h  # bottom

    if table_y0_pct >= table_y1_pct or table_y0_pct < 0:
        return None

    # Small margin (±1% of page height) for visual breathing room
    margin = 0.01
    return {
        "page": cpage,
        "x0_pct": round(x0_pct, 6),
        "x1_pct": round(x1_pct, 6),
        "y0_pct": round(max(0.0, table_y0_pct - margin), 6),
        "y1_pct": round(min(1.0, table_y1_pct + margin), 6),
        "page_w": page_w,
        "page_h": page_h,
    }


def _compute_chunk_line_ids(
    pdf_bytes: bytes, chunks: list[Chunk], line_map: dict[str, dict]
) -> list[list[str] | None]:
    """Match chunks to line_ids by text-content overlap.

    Reads per-line text from ``line_map`` entries (populated by
    ``_build_line_map``) — no need to re-open the PDF.

    Uses text-content matching instead of the original char-offset approach
    because chunk.char_start/char_end reference the markdown pipeline's
    character positions, which are a different text stream than PyMuPDF's
    ``page.get_text("text")`` output — the two never align."""
    if not line_map or not chunks:
        return [None] * len(chunks)

    # Group line_map entries by page, sorted by y position.
    # Entries with empty text are filtered out (no content to match).
    by_page: dict[int, list[tuple[str, dict]]] = {}
    for lid, info in line_map.items():
        pg = info["page"]
        if not info.get("text"):
            continue
        by_page.setdefault(pg, []).append((lid, info))
    for pg in by_page:
        by_page[pg].sort(key=lambda x: x[1]["y0_pct"])

    # Match each chunk to lines by text-content overlap
    result: list[list[str] | None] = []
    for c in chunks:
        pg = c.page
        # Table chunks: Markdown text doesn't correspond to PDF text lines.
        # Try spatial matching first (Camelot bbox → line overlap); fall back
        # to all lines on the page.
        if c.kind == "table":
            if pg >= 1 and pg in by_page:
                lids = _match_lines_spatial(c, by_page[pg])
                result.append(lids if lids else [lid for lid, _ in by_page[pg]])
            else:
                result.append(None)
            continue
        ct = (c.text or "").strip()
        if not ct or pg < 1 or pg not in by_page:
            result.append(None)
            continue
        ct_lower = ct.lower()
        lids = []
        for lid, info in by_page[pg]:
            line_text = info.get("text", "")
            if line_text and line_text.lower() in ct_lower:
                lids.append(lid)
        result.append(lids if lids else None)
    return result


def _bbox_from_line_ids(
    line_ids: list[str] | None, line_map: dict[str, dict]
) -> dict | None:
    """Compute a chunk bbox as the union of its lines' full-width bands.

    Returns {page, x0_pct, x1_pct, y0_pct, y1_pct, page_w, page_h} or None.
    Full-width bands with a small visual margin (5%-95%) look cleaner than
    0-100% and still read as "this whole region"."""
    if not line_ids or not line_map:
        return None
    lines = [line_map[lid] for lid in line_ids if lid in line_map]
    if not lines:
        return None
    l0 = lines[0]
    return {
        "page": l0["page"],
        "x0_pct": 0.05, "x1_pct": 0.95,
        "y0_pct": min(l["y0_pct"] for l in lines),
        "y1_pct": max(l["y0_pct"] + l["h_pct"] for l in lines),
        "page_w": l0["page_w"],
        "page_h": l0["page_h"],
    }


def _build_block_map(ir_document) -> dict[str, dict] | None:
    """Build a per-block geometry registry from the IR Document's blocks.

    Returns {block_id: {kind, page, x0_pct, y0_pct, x1_pct, y1_pct, page_w, page_h, text}}
    for every block that carries a bbox.  Block IDs are stable hex strings
    (b_0000, b_0001, …) — they survive re-ingestion because Docling produces
    the same items in the same order for the same PDF.

    The ``text`` field carries the block's rendered markdown snippet so the
    markdown exporter can annotate block boundaries for clickable-PDF-sync.

    Returns None when no block carries a bbox (e.g. markdown/vision path)."""
    if ir_document is None:
        return None
    blocks = getattr(ir_document, "blocks", None) or []
    reg: dict[str, dict] = {}
    for i, b in enumerate(blocks):
        bb = getattr(b, "bbox", None)
        if bb is None:
            continue
        pw = getattr(bb, "page_w", 0) or 0
        ph = getattr(bb, "page_h", 0) or 0
        if pw <= 0 or ph <= 0:
            continue
        txt = (getattr(b, "render", lambda: "")() or "").strip()
        reg[f"b_{i:04d}"] = {
            "kind": str(getattr(b, "kind", "paragraph") or "paragraph").removeprefix("BlockKind."),
            "page": int(getattr(bb, "page", 1) or 1),
            "x0_pct": round(float(bb.x0) / pw, 6) if pw else 0,
            "y0_pct": round(float(bb.y0) / ph, 6) if ph else 0,
            "x1_pct": round(float(bb.x1) / pw, 6) if pw else 0,
            "y1_pct": round(float(bb.y1) / ph, 6) if ph else 0,
            "page_w": pw,
            "page_h": ph,
            "text": txt,
        }
    return reg if reg else None


def _locate_text_span(page, text: str, pw: float, ph: float) -> dict | None:
    """Find the bounding box that covers the FULL span of `text` on `page`.
    Uses word-level extraction to match the text against page words and union
    all matched word bboxes — giving a paragraph-level box instead of a
    single-word box from search_for(text[:80]). Falls back to search_for
    on the first 80 chars when word matching produces no result."""
    # Normalize: collapse whitespace, lowercase for matching
    import re as _re
    ct_norm = _re.sub(r'\s+', ' ', (text or '').strip().lower())
    if not ct_norm:
        return None

    # Extract word-level bboxes from the page
    try:
        words = page.get_text("words")
    except Exception:  # noqa: BLE001
        words = []

    if words:
        # Build a sequence of (word_text, bbox) sorted by position (top-to-bottom, left-to-right)
        word_items = []
        for w in words:
            if len(w) >= 5 and w[4].strip():
                word_items.append({
                    "text": w[4].strip().lower(),
                    "x0": float(w[0]), "y0": float(w[1]),
                    "x1": float(w[2]), "y1": float(w[3]),
                })
        # Sort by y first, then x (reading order)
        word_items.sort(key=lambda w: (round(w["y0"]), w["x0"]))
        word_seq = [w["text"] for w in word_items]

        # ── Strategy 1: first-word / last-word matching ──
        # Match the chunk text span by locating its first word and a
        # distinctive tail word in the page word sequence, then union
        # all word bboxes between them.  Does NOT use substring find
        # (which breaks on character-level normalisation mismatches).
        needle_words = ct_norm.split()
        if len(needle_words) >= 2:
            first_word = needle_words[0]
            first_idx = None
            for wi, wt in enumerate(word_seq):
                if wt == first_word:
                    first_idx = wi
                    break
            last_idx = None
            # Walk backwards through needle words — chunk text is often
            # truncated mid-word ("regulations" → "regulatio"), so we
            # look for the first tail word that actually exists.
            for ni in range(len(needle_words) - 1, 0, -1):
                cand = needle_words[ni]
                if len(cand) <= 2:          # skip noise words
                    continue
                for wi in range(len(word_seq) - 1,
                                first_idx if first_idx is not None else 0, -1):
                    if word_seq[wi] == cand:
                        last_idx = wi
                        break
                if last_idx is not None:
                    break
            if first_idx is not None and last_idx is not None and last_idx >= first_idx:
                matched = word_items[first_idx:last_idx + 1]
                return {
                    "page": 1,  # caller overwrites with actual page
                    "x0": min(w["x0"] for w in matched),
                    "y0": min(w["y0"] for w in matched),
                    "x1": max(w["x1"] for w in matched),
                    "y1": max(w["y1"] for w in matched),
                    "page_w": pw,
                    "page_h": ph,
                }

        # ── Strategy 2: substring find → word-range bbox ──
        # Less robust (character-level differences cause misses), but
        # preserved as fallback for chunks whose first word isn't found.
        page_text = " ".join(word_seq)
        idx = page_text.find(ct_norm[:80])
        if idx < 0 and len(ct_norm) > 40:
            idx = page_text.find(ct_norm[:40])
        if idx >= 0:
            char_pos = 0
            matched = []
            for wi, wt in enumerate(word_seq):
                wlen = len(wt)
                if char_pos + wlen > idx and char_pos < idx + len(ct_norm[:200]):
                    matched.append(word_items[wi])
                char_pos += wlen + 1  # +1 for the join space
            if matched:
                return {
                    "page": 1,
                    "x0": min(w["x0"] for w in matched),
                    "y0": min(w["y0"] for w in matched),
                    "x1": max(w["x1"] for w in matched),
                    "y1": max(w["y1"] for w in matched),
                    "page_w": pw,
                    "page_h": ph,
                }

    # Fallback: search_for with the first 80 chars
    needle = (text or "").strip()[:80]
    if not needle:
        return None
    try:
        rects = page.search_for(needle, quads=False)
    except Exception:  # noqa: BLE001
        rects = []
    if not rects:
        return None
    r = rects[0]
    return {
        "page": 1,
        "x0": float(r.x0),
        "y0": float(r.y0),
        "x1": float(r.x1),
        "y1": float(r.y1),
        "page_w": pw,
        "page_h": ph,
    }


# ---- Composed pipeline ------------------------------------------------------
def ingest_document(db: Session, document_pk: int, tenant_id: str) -> dict:
    """Run the full pipeline for one document. Tenant is passed explicitly
    because the worker context bootstraps each job — we don't trust whatever
    contextvar happens to be set."""
    set_current_tenant(tenant_id)
    doc = db.scalar(
        select(Document).where(Document.tenant_id == tenant_id, Document.pk == document_pk)
    )
    if doc is None:
        raise RuntimeError(f"Document pk={document_pk} not found in tenant {tenant_id}")
    if not doc.s3_key:
        raise RuntimeError(f"Document pk={document_pk} has no s3_key — nothing to ingest")

    doc.ingestion_status = "processing"
    doc.ingestion_error = None
    db.commit()

    try:
        # 0. Document Strategist — profile the document and decide the best
        #    processing strategy BEFORE the heavy pipeline runs. This replaces
        #    the old one-size-fits-all defaults with per-document routing.
        from app.services.document_strategist import (
            strategize, profile_from_upload,
        )

        # 1. Fetch the file bytes from object storage. We use the streaming
        # helper and accumulate — typical compliance PDFs are MBs, not GBs.
        buf = io.BytesIO()
        for chunk in storage.stream_object(doc.s3_key):
            buf.write(chunk)
        raw_bytes = buf.getvalue()

        # Build the profile from upload-time metadata + cheap sampling of the bytes.
        _profile = profile_from_upload(
            raw_bytes=raw_bytes,
            mime=doc.mime_type or "",
            size_bytes=len(raw_bytes),
            name=doc.name or "",
        )
        # Blend in what we already know from the Document row (real page count
        # from upload, classifier output if already run).
        if doc.pages and doc.pages > 1:
            _profile.pages = doc.pages
        if doc.doc_type:
            _profile.doc_type = doc.doc_type
            _profile.doc_type_confidence = doc.doc_type_confidence

        strategy = strategize(_profile)
        log.info(
            "strategist: doc pk=%s name=%r → %s",
            doc.pk, doc.name, strategy.reason,
        )

        # ── Apply strategy flags ──────────────────────────────────────────
        # Override the global embed backend with the per-document choice.
        _embed_backend_override = strategy.embed_backend
        _use_contextual = strategy.use_contextual_retrieval
        _use_pdfplumber = strategy.use_pdfplumber_tables

        # 2. Parse — pick the right parser by MIME.
        #
        # Image MIMEs go through vision OCR. PDFs first try PyMuPDF; if no
        # extractable text is found (scanned PDFs / image-only PDFs), fall
        # back to per-page vision OCR. The rest of the pipeline doesn't
        # care which parser produced the (page, text) tuples.
        mime = (doc.mime_type or "").lower()
        name_lower = (doc.name or "").lower()
        is_image = mime.startswith("image/") or mime in {
            "image/heic", "image/heif", "image/avif",
        }
        is_csv = (
            mime in {"text/csv", "application/csv", "text/tab-separated-values"}
            or name_lower.endswith((".csv", ".tsv"))
        )
        # T1.2 · new format routing
        is_docx = (
            mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or name_lower.endswith(".docx")
        )
        is_xlsx = (
            mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            or name_lower.endswith(".xlsx")
        )
        is_text = (
            mime == "text/plain"
            or name_lower.endswith((".txt", ".md", ".markdown", ".log"))
        )
        is_eml = (
            mime == "message/rfc822"
            or name_lower.endswith(".eml")
        )
        is_pptx = (
            mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or name_lower.endswith(".pptx")
        )
        # Office/ODF formats LibreOffice converts → PDF (opt-in; no-op if soffice absent).
        _ext = os.path.splitext(name_lower)[1]
        is_office_convert = _ext in _OFFICE_CONVERT_EXTS
        vision_pages: set[int] = set()  # G3 · page numbers that went through OCR
        _dm = get_settings().doc_model   # Phase 2 · structured-vision when on
        ir_document = None               # Phase 4 · set for pure-native structured PDFs → block-aware chunking
        if is_image:
            pages = ingestion_vision.parse_image(
                raw_bytes, mime, db=db, tenant_id=doc.tenant_id, structured=_dm,
            )
            if _dm:
                pages = _vision_md_to_pages(pages)
            vision_pages = {p for p, _ in pages}
        elif is_csv:
            # CSV/TSV bank/CC statements etc: no vision needed. Parse to a
            # structured Markdown table (stdlib csv reader — quote/newline/
            # delimiter-aware) so retrieval and the fact extractor see real rows
            # instead of raw delimited bytes. The extractor still reads the body
            # and pulls every row into top_transactions[] via the bank_statement
            # schema — a Markdown table reads more reliably than raw CSV.
            pages = parse_csv(raw_bytes)
        elif is_docx:
            parsed = parse_docx(raw_bytes, db=db, tenant_id=doc.tenant_id, return_ir=_dm)
            if _dm and isinstance(parsed, tuple):
                pages, ir_document = parsed
            else:
                pages = parsed
        elif is_xlsx:
            pages = parse_xlsx(raw_bytes)
        elif is_text:
            pages = parse_text(raw_bytes)
        elif is_eml:
            pages = parse_eml(raw_bytes)
        elif is_pptx:
            parsed = parse_pptx(raw_bytes, db=db, tenant_id=doc.tenant_id, return_ir=_dm)
            if _dm and isinstance(parsed, tuple):
                pages, ir_document = parsed
            else:
                pages = parsed
        elif is_office_convert:
            # Convert legacy MS Office / ODF / RTF → PDF via LibreOffice, then
            # parse as a PDF. No-op (None) when soffice isn't installed (slim
            # image) → raise a clear, actionable error rather than mis-parsing.
            pdf_bytes = libreoffice_to_pdf(raw_bytes, _ext)
            if pdf_bytes is None:
                raise RuntimeError(
                    f"{_ext} requires LibreOffice (not installed in this image). "
                    "Re-save as .pdf/.docx/.xlsx/.pptx, or deploy the full image "
                    "with DOCAIQ_OFFICE_CONVERT enabled."
                )
            pages = _parse_pdf_pages(pdf_bytes)
        else:
            # G13 · Try Docling first for IR capture (layout-aware, multi-column)
            if _dm and is_enabled("documents_docling_enabled", False):
                dl = parse_pdf_docling(raw_bytes)
                if dl is not None:
                    pages, ir_document = dl
                else:
                    pages = parse_pdf(raw_bytes)  # Docling failed → fall back
            elif _dm:
                try:
                    ir_document = parse_pdf_structured(raw_bytes)
                    pages = ir_document.to_pages()
                except Exception as e:  # noqa: BLE001
                    log.warning("doc-model: structured PDF parse failed (%s); flat parse", e)
                    ir_document, pages = None, parse_pdf(raw_bytes)
            else:
                pages = parse_pdf(raw_bytes)
            # Whole-doc fallback: zero text anywhere → vision OCR every page.
            if not any((t or "").strip() for _, t in pages):
                log.info("ingest pk=%s: no text via PyMuPDF, falling back to vision OCR", doc.pk)
                ir_document = None  # vision/mixed paths chunk from serialised text (flat) for now
                pages = ingestion_vision.parse_pdf_via_vision(
                    raw_bytes, db=db, tenant_id=doc.tenant_id, structured=_dm,
                )
                if _dm:
                    pages = _vision_md_to_pages(pages)
                vision_pages = {p for p, _ in pages}
            else:
                # Partial-text fallback: some pages have text, others don't
                # (mixed text+image PDFs — common with bank/CC statements
                # whose later pages are scanned). OCR only the empty pages
                # so we don't waste vision calls on pages PyMuPDF handled.
                _TEXT_THRESHOLD = 40  # chars; less than this = effectively blank
                empty_pages = [p for p, t in pages if not t or len(t.strip()) < _TEXT_THRESHOLD]
                if empty_pages:
                    log.info(
                        "ingest pk=%s: %d page(s) image-only (%s) — vision-OCR'ing those, keeping the rest as-is",
                        doc.pk, len(empty_pages), empty_pages,
                    )
                    ocr_pages = ingestion_vision.parse_pdf_pages_via_vision(
                        raw_bytes, page_indexes_1based=empty_pages,
                        db=db, tenant_id=doc.tenant_id, structured=_dm,
                    )
                    if _dm:
                        ocr_pages = _vision_md_to_pages(ocr_pages)
                    by_page = {p: t for p, t in ocr_pages}
                    pages = [(p, by_page.get(p, t)) for p, t in pages]
                    vision_pages = set(empty_pages)
                    ir_document = None  # native+vision merge → chunk from serialised text

        # Strip NUL/control bytes from extracted text ONCE here, so every
        # downstream consumer (chunks, doc summary, extraction, entities, and the
        # Postgres chunk INSERT) sees safe text. Without this, a PDF whose text
        # carries a NUL (0x00) byte crashes ingestion at the chunk write.
        pages = [(p, _sanitize_text(t)) for p, t in pages]

        # G3 · score the OCR'd pages so low-confidence scans surface for review.
        # NULL for non-OCR docs (PyMuPDF/native text is assumed clean).
        from app import ocr_quality as _ocrq
        doc.ocr_quality = _ocrq.summarize_pages(pages, vision_pages)

        # Phase 4 · block-aware chunking straight from the IR when we have a pure
        # structured Document (native PDF, no vision merge); else the flat chunker.
        chunks = chunk_document(ir_document) if (_dm and ir_document is not None) else chunk_pages(pages)
        if not chunks:
            raise RuntimeError(
                "No extractable text via PyMuPDF or vision OCR. Check the file is readable and the OpenRouter key is configured."
            )

        # P9.5 · table-preserving extraction. For native-PDF uploads (not
        # images / CSV / docx / etc, and not vision-OCR'd), pull tables via
        # pdfplumber and append them as kind='table' Markdown chunks so the
        # grid structure survives into retrieval + the materialized markdown.
        # Appended BEFORE embedding/bbox/write so the chunks list stays aligned.
        is_pdf_native = not (is_image or is_csv or is_docx or is_xlsx or is_text
                             or is_eml or is_pptx or is_office_convert)
        # Run pdfplumber on the ORIGINAL PDF whenever the source is a native PDF —
        # even in mixed text+image mode (used_vision). pdfplumber reads the real
        # PDF, so it recovers tables from the NATIVE pages and simply finds nothing
        # on the scanned/image pages. Previously `and not used_vision` disabled ALL
        # table extraction the moment a single page was OCR'd (e.g. a bank statement
        # with one scanned page), silently losing every table on the native pages.
        if is_pdf_native and _use_pdfplumber:
            try:
                table_chunks = extract_table_chunks(raw_bytes)
                if table_chunks:
                    chunks = list(chunks) + table_chunks
            except Exception as e:  # noqa: BLE001
                log.warning("table extraction skipped (non-fatal): %s", e)

        # M43.P1 · Contextual Retrieval. Generate ~50-100 token situating
        # context per chunk via the cheap LLM cascade. Prepended to the
        # chunk text BEFORE embedding so retrieval lands in a richer
        # semantic neighborhood (+35-49% recall on diverse corpora).
        # Gated on settings.contextual_retrieval_enabled — when disabled
        # we embed the raw chunk text exactly as before (full back-compat).
        from app.config import get_settings as _gs
        from app.contextual import generate_contexts as _gen_ctx, embedding_input as _embed_in
        _settings = _gs()
        if _use_contextual and is_enabled("contextual_retrieval_enabled", True):
            doc_full_text = "\n\n".join(text_ for _, text_ in pages)
            contexts = _gen_ctx(
                doc_text=doc_full_text,
                doc_name=doc.name or doc.id_external or "document",
                chunk_texts=[c.text for c in chunks],
            )
            # embed (context + text) when available; else raw text
            vectors = []
            BATCH = 64
            embed_inputs = [_embed_in(c.text, contexts[i]) for i, c in enumerate(chunks)]
            for i in range(0, len(embed_inputs), BATCH):
                vectors.extend(embed(embed_inputs[i:i+BATCH]))
        else:
            contexts = [""] * len(chunks)
            vectors = embed_chunks(chunks)

        # Retrieval Step 2 · also compute the v2 (1024d) embedding for NEW uploads when active,
        # so they're retrievable under embed_v2 (retrieval queries the embedding_v2 column then).
        # Non-fatal: a provider hiccup leaves v2 NULL for this doc; v1 is unaffected.
        v2_vectors = None
        if _settings.embed_v2_active:
            try:
                from app.embeddings import embed_v2 as _embed_v2
                from app.contextual import embedding_input as _embed_in
                # Same contextual-retrieval representation as v1 (context + text) — so v2 keeps
                # the recall boost of Contextual Retrieval, and query↔chunk spaces stay aligned.
                # Degrades to raw text when there's no context (_embed_in fallback).
                _v2_src = [_embed_in(c.text, contexts[i]) for i, c in enumerate(chunks)]
                v2_vectors = []
                for i in range(0, len(_v2_src), 32):
                    v2_vectors.extend(_embed_v2(_v2_src[i:i + 32], backend=_embed_backend_override))
            except Exception as e:  # noqa: BLE001
                log.warning("ingest pk=%s · embed_v2 skipped (non-fatal): %s", doc.pk, e)
                v2_vectors = None

        # M11.7: compute bounding box per chunk via PyMuPDF page.search_for.
        # Gate on "is the source a native PDF" — NOT on used_vision. For a native
        # PDF that went through mixed-mode OCR, search_for locates chunks on the
        # native pages (real bboxes → clickable citations) and naturally returns
        # None for chunks whose text came from OCR'ing an image page (that text
        # isn't in the PDF). For image uploads / office conversions raw_bytes isn't
        # a searchable PDF, so we skip and the chat panel shows the whole page.
        # Extract word-level text layer for reverse bbox lookup
        if mime == "application/pdf" or name_lower.endswith(".pdf"):
            _tl = None
            # PyMuPDF word extraction
            if not _tl:
                try:
                    _tl_buf = io.BytesIO()
                    for chunk in storage.stream_object(doc.s3_key):
                        _tl_buf.write(chunk)
                    _tl = _extract_text_layer(_tl_buf.getvalue(), (doc.pages or 0) or 999)
                except Exception:
                    pass
            if _tl:
                if doc.extracted_fields is None:
                    doc.extracted_fields = {}
                doc.extracted_fields["text_layer"] = _tl
                from sqlalchemy.orm.attributes import flag_modified
                flag_modified(doc, "extracted_fields")

        # ── Line-ID bbox pipeline ─────────────────────────────────────────
        # Build a per-line geometry map from the native PDF, assign each chunk
        # the line_ids it spans via char-position overlap, and compute bboxes
        # as the union of full-width line bands.  Falls back to the existing
        # _bboxes_for_chunks word-matching for non-PDF or pre-0106 docs.
        line_map = {}
        chunk_line_ids: list[list[str] | None] = [None] * len(chunks)
        if is_pdf_native:
            line_map = _build_line_map(raw_bytes)
            if line_map:
                doc.line_map = line_map
                from sqlalchemy.orm.attributes import flag_modified
                flag_modified(doc, "line_map")
                chunk_line_ids = _compute_chunk_line_ids(raw_bytes, chunks, line_map)
                log.info("ingest pk=%s · line_map: %d lines across %d pages",
                         doc.pk, len(line_map), len({v["page"] for v in line_map.values()}))

        if _dm and ir_document is not None:
            # Phase 5 · start with IR block bboxes (already computed during
            # chunking — free).  Try the cheap line-based pipeline next for gaps;
            # only fall back to expensive word-matching when both are unavailable.
            bboxes = [c.bbox for c in chunks]  # IR block bboxes come first

            # ── Block registry ───────────────────────────────────────────
            # Persist per-block geometry from the parser IR so the three-pane
            # sync and diagnostics can reference stable block identities.
            # Chunk bboxes already flow through c.bbox (union of composing
            # blocks) — this map is the source-of-truth for per-block precision.
            _bm = _build_block_map(ir_document)
            if _bm:
                doc.block_map = _bm
                from sqlalchemy.orm.attributes import flag_modified
                flag_modified(doc, "block_map")
                log.info("ingest pk=%s · block_map: %d blocks with bbox",
                         doc.pk, len(_bm))
        elif is_pdf_native:
            bboxes = [c.bbox for c in chunks]
        else:
            bboxes = [None] * len(chunks)

        # ── Line-ID bbox gap-fill ───────────────────────────────────────
        # Full-width line-band unions are cheaper than word-matching (dict
        # lookups vs. text search) and more reliable (no whitespace mismatches,
        # no x-coordinate ambiguity).  Fill as many gaps as possible here first.
        if line_map:
            _line_bboxes = 0
            for i, lids in enumerate(chunk_line_ids):
                if bboxes[i] is not None:
                    continue
                _c = chunks[i]
                tb = _bbox_from_table_bbox(_c, line_map) if _c.table_bbox else None
                if tb:
                    bboxes[i] = tb
                    _line_bboxes += 1
                else:
                    lb = _bbox_from_line_ids(lids, line_map)
                    if lb:
                        bboxes[i] = lb
                        _line_bboxes += 1
            if _line_bboxes:
                log.info("ingest pk=%s · line-based bboxes: %d/%d chunks",
                         doc.pk, _line_bboxes, len(chunks))

            # ── Word-matching gap-fill (expensive) ─────────────────────
            # Only run on chunks that still have no bbox — line-based fill
            # handles the majority case so this runs on a small subset.
            _gaps = [i for i, bb in enumerate(bboxes) if bb is None]
            if _gaps and is_pdf_native:
                _gap_chunks = [chunks[i] for i in _gaps]
                _gap_bboxes = _bboxes_for_chunks(raw_bytes, _gap_chunks)
                for j, idx in enumerate(_gaps):
                    if _gap_bboxes[j] is not None:
                        bboxes[idx] = _gap_bboxes[j]
        elif is_pdf_native and (_dm is None or ir_document is None):
            # No line_map — fall back to full word-matching for every chunk
            bboxes = _bboxes_for_chunks(raw_bytes, chunks)

        # M44.P11.2 · PII-at-rest. Tokenize the text we STORE in our own DB
        # (chunk rows → and, downstream, the materialized artifacts the
        # materializer builds from them) so persisted content carries
        # placeholders like [CREDIT_CARD_1]; real values are encrypted in
        # pii_vault. Embeddings, entities and bboxes above were computed on the
        # ORIGINAL text, so retrieval + graph quality are unaffected. New
        # uploads only — existing docs are untouched until re-ingested.
        chunk_texts = [c.text for c in chunks]
        if _settings.pii_protect_at_rest:
            from app import pii_vault
            try:
                pii_vault.tokenize_document(
                    db, tenant_id=tenant_id, document_pk=doc.pk,
                    text="\n\n".join(chunk_texts),
                )
                v2t = pii_vault.value_to_token_map(db, doc.pk)
                if v2t:
                    chunk_texts = [pii_vault.apply_mapping(t, v2t) for t in chunk_texts]
                    doc.pii_protected = True
                    doc.pii_revealed = False
                    log.info("ingest pk=%s · PII-at-rest: tokenized %d value(s)", doc.pk, len(v2t))
            except Exception as e:  # noqa: BLE001 — never block ingest on this
                log.warning("PII-at-rest tokenization skipped (non-fatal): %s", e)

        # 3. Replace any prior chunks + entities for this doc (re-ingest is idempotent).
        db.execute(delete(Entity).where(Entity.document_pk == doc.pk))
        db.execute(delete(DocumentChunk).where(DocumentChunk.document_pk == doc.pk))
        chunk_rows: list[DocumentChunk] = []
        for i, (c, v, bbox) in enumerate(zip(chunks, vectors, bboxes)):
            row = DocumentChunk(
                tenant_id=tenant_id,
                document_pk=doc.pk,
                chunk_index=i,
                text=chunk_texts[i],  # M44.P11.2 · tokenized when pii_protect_at_rest
                page=c.page,
                char_start=c.char_start,
                char_end=c.char_end,
                kind=c.kind,  # P9.5 · 'text' or 'table'
                embedding=v,
                embedding_v2=(v2_vectors[i] if v2_vectors is not None else None),
                bbox=bbox,
                line_ids=(chunk_line_ids[i] if i < len(chunk_line_ids) else None),
                block_ids=list(c.block_ids) if c.block_ids else None,
                # M43.P1 · empty string when contextual retrieval is off
                # or the LLM call failed for this chunk. The retrieval
                # layer treats empty == "no context available".
                context_summary=_sanitize_text(contexts[i]) or None,  # LLM-sourced → strip defensively
                pipeline_version=get_settings().pipeline_version,
            )
            db.add(row)
            chunk_rows.append(row)
        db.flush()  # populate chunk_rows[*].pk so entities can FK to them

        # 4. Entity pass. Regex (deterministic, compliance vocab) and/or LLM NER
        #    (general free-text entities: people/orgs/locations/products/clauses)
        #    per DOCAIQ_NER_BACKEND. "both" keeps regex precision + LLM coverage.
        ner_backend = getattr(get_settings(), "ner_backend", "regex")
        entity_count = 0
        if ner_backend in ("regex", "both"):
            for c, row in zip(chunks, chunk_rows):
                for ent in extract_entities(c.text):
                    db.add(Entity(
                        tenant_id=tenant_id,
                        document_pk=doc.pk,
                        chunk_pk=row.pk,
                        kind=ent.kind,
                        text=ent.text,
                        page=c.page,
                        entity_metadata=ent.metadata,
                    ))
                    entity_count += 1
        if ner_backend in ("llm", "both"):
            # One cheap LLM call over the doc's chunks; writes source='llm_ner'
            # rows in their own GraphRun (auditable/idempotent). Never fatal to
            # ingestion — a provider/parse failure just yields 0 NER entities.
            try:
                from app.agents import ner_extractor
                entity_count += ner_extractor.run(
                    db, doc, [(c.page, c.text) for c in chunks]
                )
            except Exception as e:  # noqa: BLE001
                log.warning("ingest: LLM NER pass failed for doc pk=%s: %s", doc.pk, e)

        # 5. Update document metadata — real page count, ready status.
        doc.pages = max(p for p, _ in pages) if pages else 1

        # M47 · free-plan page cap — enforced at the point the TRUE page count is
        # known, so every ingestion path (upload, Drive sync/import, link-pull) is
        # covered, not just direct upload. An over-limit doc is failed with a clear
        # message rather than silently consuming a slot.
        try:
            from app.services import subscriptions as _subs
            _cap = _subs.page_cap_for(db, owner_user_id=doc.owner_user_id) if doc.owner_user_id else None
        except Exception:  # noqa: BLE001 — never break ingestion on a plan-config read
            _cap = None
        if _cap is not None and (doc.pages or 0) > _cap:
            _limit_txt = "single-page documents only" if _cap == 1 else f"up to {_cap} pages per document"
            doc.ingestion_status = "failed"
            doc.ingestion_error = (
                f"This document has {doc.pages} pages — the free plan allows {_limit_txt}. "
                f"Split it into single-page files, or upgrade for multi-page documents."
            )
            db.commit()
            log.info("ingest: doc pk=%s over free page cap (%d > %d) — marked failed",
                     doc.pk, doc.pages, _cap)
            return {"document_pk": doc.pk, "pages": doc.pages, "chunks": len(chunks),
                    "entities": entity_count, "failed": "page_cap"}
        doc.ingestion_status = "ready"

        # M47 · Indexing quality critic — fast check always runs (free).
        # LLM critique runs when documents_indexing_critic=true.
        try:
            from app.agents.indexing_critic import run_fast, run as critic_run
            _chunk_texts = [c.text for c in chunks if c.text]
            _fast = run_fast(doc.name or "", doc.doc_type or "", _chunk_texts)
            for issue in _fast.get("issues", []):
                log.info("indexing_critic: %s · %s", doc.name, issue)
            if is_enabled("documents_indexing_critic", False):
                _entities = [{"kind": e.kind, "text": e.text, "canonical": e.canonical}
                             for e in (doc.entities or [])[:30]]
                critic_run(db, doc.name or "", doc.doc_type or "",
                           _chunk_texts, _entities, tenant_id=tenant_id)
                # Result logged inside critic_run()
        except Exception:
            pass  # never block ingestion

        # M46 · retention. The original is KEPT here. Purging the blob the instant
        # ingestion (chunking) finished was wrong on three counts: the document
        # viewer needs the original to render, downstream vision/KYC extraction
        # runs in LATER tasks and needs it too, and re-open broke. Time-based
        # retention (e.g. purge connector originals 24h after the full pipeline is
        # ready, with re-pull-on-view as the fallback) belongs in a separate
        # scheduled cleanup, not an eager mid-pipeline delete. retain_original
        # still governs that future job; nothing is purged inline anymore.

        db.commit()
        log.info(
            "Ingested document pk=%s tenant=%s pages=%d chunks=%d entities=%d",
            doc.pk, tenant_id, doc.pages, len(chunks), entity_count,
        )
        return {"document_pk": doc.pk, "pages": doc.pages, "chunks": len(chunks), "entities": entity_count}
    except Exception as e:
        db.rollback()
        # Re-fetch in a fresh transaction so we can persist the error.
        fresh = db.scalar(select(Document).where(Document.pk == document_pk))
        if fresh is not None:
            fresh.ingestion_status = "failed"
            fresh.ingestion_error = str(e)[:1000]
            db.commit()
        log.exception("Ingestion failed for document pk=%s", document_pk)
        raise
